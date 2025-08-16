from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import List, Tuple

FONT_NAME = "Times New Roman"
FONT_SIZE_PT = 20

# Helper functions

def set_times_20_textframe(tf):
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(FONT_SIZE_PT)


def find_slide_indices_by_titles(prs: Presentation, titles: List[str]) -> List[int]:
    idx = []
    for i, s in enumerate(prs.slides):
        title = s.shapes.title.text.strip() if s.shapes.title and s.shapes.title.has_text_frame else ""
        if title in titles:
            idx.append(i)
    return idx


def find_slide_index_by_title(prs: Presentation, title: str) -> int:
    for i, s in enumerate(prs.slides):
        t = s.shapes.title.text.strip() if s.shapes.title and s.shapes.title.has_text_frame else ""
        if t == title:
            return i
    return -1


def delete_slides(prs: Presentation, indices: List[int]):
    # Must delete from highest to lowest
    for i in sorted(indices, reverse=True):
        xml_slides = prs.slides._sldIdLst  # pylint: disable=protected-access
        slide_id = xml_slides[i].rId
        prs.part.drop_rel(slide_id)
        xml_slides.remove(xml_slides[i])


def add_title_content_slide(prs: Presentation, title: str, bullets: List[Tuple[int, str]]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    ph = slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    first = True
    for lvl, text in bullets:
        if first:
            tf.text = text
            tf.paragraphs[0].level = lvl
            first = False
        else:
            p = tf.add_paragraph()
            p.text = text
            p.level = lvl
    set_times_20_textframe(tf)
    return slide


def add_table(slide, rows: int, cols: int, left_in: float, top_in: float, width_in: float, height_in: float, headers: List[str], data: List[List[str]]):
    left, top = Inches(left_in), Inches(top_in)
    width, height = Inches(width_in), Inches(height_in)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    # headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT_NAME
                r.font.size = Pt(FONT_SIZE_PT)
    # data
    for i in range(1, rows):
        for j in range(cols):
            val = data[i-1][j] if i-1 < len(data) and j < len(data[i-1]) else ""
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT_NAME
                    r.font.size = Pt(FONT_SIZE_PT)
    return table_shape


if __name__ == "__main__":
    base = Path("lession_expanded2.pptx")
    if not base.exists():
        base = Path("lession_expanded.pptx")
    if not base.exists():
        base = Path("lession.pptx")
    prs = Presentation(str(base))

    # Merge: 3GPP events D2 + A4/A5 into one fuller slide
    d2_title = "Event D2 Mechanics (3GPP TS 38.331)"
    a4a5_title = "A4/A5 Practical Tuning (3GPP TS 38.331)"
    d2_idx = find_slide_index_by_title(prs, d2_title)
    a4_idx = find_slide_index_by_title(prs, a4a5_title)
    if d2_idx >= 0 and a4_idx >= 0:
        insert_pos = min(d2_idx, a4_idx)  # where merged slide will be inserted
        slide = add_title_content_slide(
            prs,
            "3GPP Event Mechanics & Tuning",
            [
                (0, "D2 Enter: Ml1 – Hys > Thresh1 AND Ml2 + Hys < Thresh2"),
                (1, "Leave: Ml1 + Hys < Thresh1 OR Ml2 – Hys > Thresh2"),
                (0, "A4 Enter/Leave via Hys with thresholds; suppress ping‑pong"),
                (1, "A5 dual thresholds: Mp < Thresh1 AND Mn > Thresh2"),
                (0, "Tune Hys + TTT + Ofn/Ocn per Doppler & speed")
            ]
        )
        # Table with definitions
        headers = ["Param", "Meaning"]
        data = [
            ["Ml1/Ml2", "UE distance vs moving ref (SIB19 / configured)"],
            ["Hys", "Hysteresis to suppress ping‑pong"],
            ["Thresh", "Threshold(s) for enter/leave"],
            ["Ofn/Ocn", "Freq/cell offsets to steer preference"],
            ["Mp/Mn", "Serving/neighbor measurement"]
        ]
        add_table(slide, rows=len(data)+1, cols=2, left_in=0.8, top_in=4.0, width_in=8.8, height_in=2.5, headers=headers, data=data)
        # Delete originals
        delete_slides(prs, [d2_idx, a4_idx])

    # Merge: DRL + MARL + Nash-SAC
    t_drl = "Protocol‑Learning DRL for Handover"
    t_marl = "Multi‑Agent RL & Stability"
    t_nash = "Game‑Theoretic RL: Nash‑SAC"
    i_drl = find_slide_index_by_title(prs, t_drl)
    i_marl = find_slide_index_by_title(prs, t_marl)
    i_nash = find_slide_index_by_title(prs, t_nash)
    if min(i for i in [i_drl, i_marl, i_nash] if i >= 0) >= 0:
        insert_pos = min(i for i in [i_drl, i_marl, i_nash] if i >= 0)
        slide = add_title_content_slide(
            prs,
            "AI for Handover: DRL • MARL • Nash‑SAC",
            [
                (0, "DRL: End‑to‑end policy; state→trigger; fewer outages"),
                (0, "MARL: Distributed agents; load stability via rewards"),
                (0, "Nash‑SAC: Game‑theoretic stability under high load")
            ]
        )
        headers = ["Method", "Key Idea", "Reported Gains"]
        data = [
            ["Protocol‑Learning DRL", "Learn HO rules from data", "Lower outage vs fixed rules"],
            ["MARL", "Local info; implicit coordination", "Load balance; fewer oscillations"],
            ["Nash‑SAC", "Nash‑consistent SAC policies", "Stable under multi‑user load"]
        ]
        add_table(slide, rows=len(data)+1, cols=3, left_in=0.5, top_in=3.8, width_in=9.2, height_in=3.0, headers=headers, data=data)
        delete_slides(prs, [i for i in [i_drl, i_marl, i_nash] if i >= 0])

    # Merge: Transformer + CHO‑XGBoost
    t_tx = "Transformer‑based Channel Prediction"
    t_xgb = "CHO Trigger Prediction (XGBoost)"
    i_tx = find_slide_index_by_title(prs, t_tx)
    i_xgb = find_slide_index_by_title(prs, t_xgb)
    if i_tx >= 0 and i_xgb >= 0:
        insert_pos = min(i_tx, i_xgb)
        slide = add_title_content_slide(
            prs,
            "Prediction for Proactive & CHO",
            [
                (0, "Transformer: long‑range attention for channel"),
                (1, "Inputs: RSRP history, Doppler, geometry"),
                (0, "XGBoost: supervised CHO success prediction"),
                (1, "Features: geometry, timers, KPIs; lower failed HOs")
            ]
        )
        headers = ["Model", "Inputs", "Outcome"]
        data = [
            ["Transformer", "RSRP, Doppler, positions", "10–30s ahead channel"],
            ["XGBoost", "Geometry, timing, KPIs", "CHO trigger success/prob."]
        ]
        add_table(slide, rows=len(data)+1, cols=3, left_in=0.5, top_in=4.0, width_in=9.2, height_in=2.8, headers=headers, data=data)
        delete_slides(prs, [i_tx, i_xgb])

    # Expand: Measurement Latency & Doppler Mitigation
    t_lat = "Measurement Latency & Doppler Mitigation"
    i_lat = find_slide_index_by_title(prs, t_lat)
    if i_lat >= 0:
        s = prs.slides[i_lat]
        ph = s.placeholders[1]
        tf = ph.text_frame
        for extra in [
            (0, "Mitigate latency: TTT tuning; shorter report intervals"),
            (0, "RSRP smoothing: EWMA/median filters"),
            (0, "Doppler‑aware frequency correction in PHY"),
            (0, "Fallback to D2/CHO when KPIs degrade")
        ]:
            p = tf.add_paragraph()
            p.text = extra[1]
            p.level = extra[0]
        set_times_20_textframe(tf)
        headers = ["Problem", "Mitigation"]
        data = [
            ["Latency", "TTT decrease; faster reporting"],
            ["Doppler", "Compensation; pre‑comp in link est."],
            ["Noise", "Filtering; robust thresholds"],
            ["Outages", "CHO/D2 fallback policies"]
        ]
        add_table(s, rows=len(data)+1, cols=2, left_in=0.6, top_in=4.5, width_in=9.0, height_in=2.5, headers=headers, data=data)

    # Expand: Graph‑based Planning Details
    t_graph = "Graph‑based Planning Details"
    i_graph = find_slide_index_by_title(prs, t_graph)
    if i_graph >= 0:
        s = prs.slides[i_graph]
        tf = s.placeholders[1].text_frame
        for extra in [
            (0, "Heuristics: beam pruning; time windowing"),
            (0, "Constraints: HO budget; outage penalties"),
            (0, "Backhaul‑aware edge weights")
        ]:
            p = tf.add_paragraph(); p.text = extra[1]; p.level = extra[0]
        set_times_20_textframe(tf)
        headers = ["Aspect", "Notes"]
        data = [
            ["Nodes", "(sat, beam, time)"],
            ["Edges", "quality + HO cost + backhaul"],
            ["Pruning", "visibility windows; k‑best"],
        ]
        add_table(s, rows=len(data)+1, cols=2, left_in=0.6, top_in=4.2, width_in=9.0, height_in=2.5, headers=headers, data=data)

    # Expand: Location‑based HO with Particle Filter
    t_pf = "Location‑based HO with Particle Filter"
    i_pf = find_slide_index_by_title(prs, t_pf)
    if i_pf >= 0:
        s = prs.slides[i_pf]
        tf = s.placeholders[1].text_frame
        for extra in [
            (0, "State: position/velocity; Obs: ranges/angles"),
            (0, "Update: predict→weight→resample"),
            (0, "Trigger: geometry threshold + hysteresis")
        ]:
            p = tf.add_paragraph(); p.text = extra[1]; p.level = extra[0]
        set_times_20_textframe(tf)
        headers = ["Item", "Detail"]
        data = [
            ["Inputs", "ephemeris, UE motion"],
            ["Estimator", "particle filter"],
            ["Output", "trigger timing"],
        ]
        add_table(s, rows=len(data)+1, cols=2, left_in=0.6, top_in=4.2, width_in=9.0, height_in=2.5, headers=headers, data=data)

    # Expand: Evaluation Methodology
    t_eval = "Evaluation Methodology"
    i_eval = find_slide_index_by_title(prs, t_eval)
    if i_eval >= 0:
        s = prs.slides[i_eval]
        tf = s.placeholders[1].text_frame
        for extra in [
            (0, "Traffic: eMBB/URLLC mix; mobility profiles"),
            (0, "Simulation: time step, channel model, noise"),
            (0, "Training: episodes, exploration policy")
        ]:
            p = tf.add_paragraph(); p.text = extra[1]; p.level = extra[0]
        set_times_20_textframe(tf)
        headers = ["KPI", "Definition"]
        data = [
            ["HO rate", "#HOs per minute"],
            ["RLF", "radio link failures %"],
            ["Throughput", "Mbps per UE"],
            ["Latency", "E2E ms"],
            ["QoS%", "% time meeting QoS"]
        ]
        add_table(s, rows=len(data)+1, cols=2, left_in=0.6, top_in=4.0, width_in=9.0, height_in=2.8, headers=headers, data=data)

    # Expand: Implementation Guidelines
    t_impl = "Implementation Guidelines"
    i_impl = find_slide_index_by_title(prs, t_impl)
    if i_impl >= 0:
        s = prs.slides[i_impl]
        tf = s.placeholders[1].text_frame
        for extra in [
            (0, "Guardrails: min dwell time; max HO per minute"),
            (0, "Fallback: revert to D2/CHO on KPI drop"),
            (0, "Rollout: A/B test per region; feature flags")
        ]:
            p = tf.add_paragraph(); p.text = extra[1]; p.level = extra[0]
        set_times_20_textframe(tf)

    # Expand: ISL/Backhaul Constraints & Mobility
    t_isl = "ISL/Backhaul Constraints & Mobility"
    i_isl = find_slide_index_by_title(prs, t_isl)
    if i_isl >= 0:
        s = prs.slides[i_isl]
        tf = s.placeholders[1].text_frame
        for extra in [
            (0, "Metrics: ISL changes/hour; backhaul capacity"),
            (0, "Policy: HO with backhaul headroom criteria"),
            (0, "Forecast: predict ISL reconfig windows")
        ]:
            p = tf.add_paragraph(); p.text = extra[1]; p.level = extra[0]
        set_times_20_textframe(tf)

    out = "lession_final.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides")

