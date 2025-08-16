from pptx import Presentation
from pptx.util import Pt, Inches
from pathlib import Path
from typing import List, Tuple

FONT_NAME = "Times New Roman"
FONT_SIZE_PT = 20

SLIDES_TO_INSERT = [
    {
        "after_title": "TLE Data Processing Pipeline",
        "title": "Phase‑1 Targets & Time Window",
        "bullets": [
            "Scale targets: <2 min load; <100 m accuracy; <2 GB memory; ≥90% success",
            "Observer: NTPU (24.9441667°N, 121.3713889°E, 50 m)",
            "Time grid (code): 30‑s steps; default 200 min → 400 points",
            "Deck example: 6 h window → 720 points (configurable)",
            "Batching by 100 satellites; async gather per batch"
        ],
        "table": {
            "headers": ["Param", "Value (Code)", "Notes"],
            "rows": [
                ["time_range_minutes", "200", "Default; configurable"],
                ["time_resolution_seconds", "30", "Fixed step"],
                ["batch_size", "100", "Compute batches"],
                ["observer", "NTPU Topos", "lat/lon/alt from code"],
                ["targets", "<2m; <100m; <2GB; ≥90%", "Phase‑1 goals"],
            ]
        }
    },
    {
        "after_title": "TLE Parsing Process",
        "title": "TLE Field Extraction Map",
        "bullets": [
            "Line 1: sat no., epoch year/day → datetime (UTC)",
            "Line 2: inclination, RAAN, e, ω, M, n (rev/day)",
            "Derived: a from n; period (min); apogee/perigee altitudes",
            "Consistency: validator checks; sample_limits optional",
            "Constellation tag from source path (starlink/oneweb)"
        ],
        "table": {
            "headers": ["Field", "Slice", "Meaning"],
            "rows": [
                ["satellite_number", "L1[2:7]", "NORAD ID"],
                ["epoch (yy, ddd.ddd)", "L1[18:20], L1[20:32]", "Year/day → datetime"],
                ["inclination_deg", "L2[8:16]", "Orbit tilt"],
                ["raan_deg", "L2[17:25]", "Ascending node"],
                ["eccentricity", "L2[26:33]", "0.xxxxxxx"],
                ["arg_perigee_deg", "L2[34:42]", "ω"],
                ["mean_anomaly_deg", "L2[43:51]", "M"],
                ["mean_motion", "L2[52:63]", "n (rev/day)"]
            ]
        }
    },
    {
        "after_title": "SGP4 Orbital Propagation",
        "title": "Skyfield vs. Enhanced SGP4",
        "bullets": [
            "Primary: Skyfield EarthSatellite.at(t); subpoint(); altaz()",
            "Fallback: enhanced SGP4 (Kepler + J2; numeric E solver)",
            "Observer‑relative: elevation/azimuth/distance, range‑rate",
            "Velocity: numeric diff at Δt=1 s or 0.1 s",
            "Accuracy estimate capped at 200 m"
        ],
        "table": {
            "headers": ["Path", "Inputs", "Outputs"],
            "rows": [
                ["Skyfield", "TLE L1/L2, ts, Topos", "lat/lon/alt; el/az/dist; v (km/s)"],
                ["Enhanced SGP4", "SGP4 params, dt", "ECEF pos/vel; geo; el/az/dist"]
            ]
        }
    },
    {
        "after_title": "Coordinate System Transformations",
        "title": "Observer‑Relative & Frames",
        "bullets": [
            "Frames: ECI → ECEF → geographic; GST for Earth rotation",
            "Observer ECEF from lat/lon/alt + GST",
            "Relative vector → unit vector; elevation/azimuth from components",
            "Range‑rate = v·r̂; positive receding",
            "All timestamps normalized to UTC"
        ],
        "table": {
            "headers": ["Quantity", "Symbol", "Unit"],
            "rows": [
                ["Elevation", "el", "deg"],
                ["Azimuth", "az", "deg"],
                ["Distance", "|r|", "km"],
                ["Range‑rate", "ṙ", "km/s"],
                ["Velocity", "|v|", "km/s"]
            ]
        }
    },
    {
        "after_title": "Visibility Calculations",
        "title": "Visibility & Filtering",
        "bullets": [
            "Input: observer Topos + satellite at(t)",
            "Compute: el/az/dist; visibility if el ≥ threshold",
            "Typical threshold: 10–15° (configurable)",
            "Candidate reduction before later stages",
            "Output: boolean + geometric details"
        ],
        "table": {
            "headers": ["Threshold", "Effect"],
            "rows": [
                ["el ≥ 10°", "Broader set; more handover options"],
                ["el ≥ 15°", "Better link budget; fewer candidates"]
            ]
        }
    },
    {
        "after_title": "Output Data Structure",
        "title": "SatellitePosition Schema",
        "bullets": [
            "Time‑series per satellite stored in memory",
            "Fields include geographic + observer‑relative",
            "Velocity magnitude from numeric diff",
            "JSON export for downstream use (optional)",
            "Load/calculation stats exported on request"
        ],
        "table": {
            "headers": ["Field", "Type/Unit"],
            "rows": [
                ["timestamp", "datetime (UTC)"],
                ["latitude_deg, longitude_deg", "deg"],
                ["altitude_km", "km"],
                ["elevation_deg, azimuth_deg", "deg"],
                ["distance_km", "km"],
                ["velocity_km_s", "km/s"]
            ]
        }
    },
    {
        "after_title": "Performance & Scalability Metrics",
        "title": "Batching & Phase‑1 Targets",
        "bullets": [
            "Local TLE discovery by newest file per constellation",
            "Parallel load via asyncio; cache validity checks",
            "Compute in batches of 100; gather with exceptions",
            "Statistics exported to JSON (timestamped)",
            "Sample limits supported during development"
        ],
        "table": {
            "headers": ["Metric", "Target"],
            "rows": [
                ["Load time", "< 120 s"],
                ["Accuracy", "< 100 m"],
                ["Memory", "< 2 GB"],
                ["Success rate", "≥ 90%"],
                ["Time points", "200 min @30 s → 400 points"]
            ]
        }
    },
    {
        "after_title": "Professional Tools & Libraries",
        "title": "Local TLE Sources (Auto‑Discover)",
        "bullets": [
            "Base: netstack/tle_data/{starlink|oneweb}/tle/",
            "Pattern: starlink_*.tle; oneweb_*.tle",
            "Pick latest by mtime per constellation",
            "Fallback: synthetic dataset if none found",
            "Validator: enhanced checks for robustness"
        ],
        "table": {
            "headers": ["Constellation", "Path Pattern"],
            "rows": [
                ["Starlink", ".../starlink/tle/starlink_*.tle"],
                ["OneWeb", ".../oneweb/tle/oneweb_*.tle"],
            ]
        }
    }
]


def set_times_20(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT_NAME
            r.font.size = Pt(FONT_SIZE_PT)


def insert_slide_after(prs: Presentation, index_after: int, title: str, bullets: List[str], table_spec=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    ph = slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    first = True
    for b in bullets:
        if first:
            tf.text = b
            first = False
        else:
            para = tf.add_paragraph()
            para.text = b
    set_times_20(tf)
    # Add table if requested
    if table_spec:
        headers = table_spec.get("headers", [])
        rows = table_spec.get("rows", [])
        # place table lower to fill slide
        left, top, width, height = Inches(0.6), Inches(3.8), Inches(9.4), Inches(3.1)
        tbl_shape = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height)
        tbl = tbl_shape.table
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            set_times_20(cell.text_frame)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                set_times_20(cell.text_frame)
    # Reorder to position after index_after
    sldIdLst = prs.slides._sldIdLst  # type: ignore
    new_sldId = sldIdLst[-1]
    sldIdLst.remove(new_sldId)
    sldIdLst.insert(index_after + 1, new_sldId)


def find_index_contains(prs: Presentation, title_substr: str) -> int:
    for i, s in enumerate(prs.slides):
        t = s.shapes.title.text.strip() if s.shapes.title and s.shapes.title.has_text_frame else ""
        if title_substr.lower() in t.lower():
            return i
    return -1


if __name__ == "__main__":
    base = Path("tle.pptx")
    prs = Presentation(str(base))
    planned = []
    for spec in SLIDES_TO_INSERT:
        idx = find_index_contains(prs, spec["after_title"])
        if idx >= 0:
            planned.append((idx, spec))
    # Insert in reverse index order to keep positions stable
    planned.sort(key=lambda x: x[0], reverse=True)
    for idx, spec in planned:
        insert_slide_after(prs, idx, spec["title"], spec["bullets"], spec.get("table"))
    out = "tle_expanded.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides")

