from pptx import Presentation
from pathlib import Path

def ppt_to_markdown(pptx_path: str, md_path: str) -> None:
    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ''
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        if not title:
            title = f"Slide {i}"
        lines.append(f"# {title}")
        for shape in slide.shapes:
            if shape.has_text_frame and (not (slide.shapes.title and shape == slide.shapes.title)):
                for para in shape.text_frame.paragraphs:
                    txt = ''.join(run.text for run in para.runs).strip()
                    if not txt:
                        continue
                    indent = '  ' * getattr(para, 'level', 0)
                    lines.append(f"{indent}- {txt}")
        lines.append("")
    Path(md_path).write_text('\n'.join(lines), encoding='utf-8')

if __name__ == "__main__":
    pptx = "lession.pptx"
    md = "lession.md"
    ppt_to_markdown(pptx, md)
    print(f"Wrote {md}")

