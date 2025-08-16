from pptx import Presentation
from pptx.util import Pt
from pathlib import Path

FONT_NAME = "Times New Roman"
FONT_SIZE_PT = 20

SLIDES_TO_INSERT = [
    {
        "after_title": "3GPP Measurement Events for Handover Trigger",
        "title": "Event D2 Mechanics (3GPP TS 38.331)",
        "bullets": [
            "Moving references + ephemeris (SIB19)",
            "Enter if: Ml1 – Hys > Thresh1 AND Ml2 + Hys < Thresh2",
            "Leave if: Ml1 + Hys < Thresh1 OR Ml2 – Hys > Thresh2",
            "Units: distances (meters); Hys/Thresh in same units",
            "Source: 3GPP TS 38.331 (Rel-18)"
        ]
    },
    {
        "after_title": "3GPP Measurement Events for Handover Trigger",
        "title": "A4/A5 Practical Tuning (3GPP TS 38.331)",
        "bullets": [
            "A4 enter/leave: hysteresis (Hys) to suppress ping‑pong",
            "A5 dual threshold: Mp < Thresh1 AND Mn > Thresh2",
            "Offsets: Ofn (freq), Ocn (cell) steer preference",
            "Use TTT + Hys + offsets; calibrate by Doppler/UE speed",
            "Source: 3GPP TS 38.331 (Rel-18)"
        ]
    },
    {
        "after_title": "Measurement Report-based Handover (Traditional Approaches)",
        "title": "Measurement Latency & Doppler Mitigation",
        "bullets": [
            "Latency: 150–300 ms round trip → stale reports",
            "Doppler: rapid frequency shifts bias RSRP/RSRQ",
            "Mitigate: higher report rate, filtering, Doppler-compensation",
            "Fallback: D2/CHO when measurements are unreliable",
            "Refs: NTN surveys (2024), 3GPP NTN overview"
        ]
    },
    {
        "after_title": "Measurement Report-based Handover (AI-Enhanced Approaches)",
        "title": "Protocol‑Learning DRL for Handover",
        "bullets": [
            "End‑to‑end learning of HO rules from experience",
            "State: channel stats, geometry, timers; Action: HO trigger",
            "Reward: throughput − HO penalty − outage cost",
            "Fewer outages vs fixed rules in simulations",
            "Source: Handover Protocol Learning for LEO (arXiv:2310.20215)"
        ]
    },
    {
        "after_title": "Measurement Report-based Handover (AI-Enhanced Approaches)",
        "title": "Multi‑Agent RL & Stability",
        "bullets": [
            "Distributed agents (UEs) with local observations",
            "Coordination via environment; avoid oscillations",
            "Load‑aware MARL (GLOBECOM 2020) improves balance",
            "Design: reward shaping + entropy regularization",
            "Refs: MARL for LEO HO (survey, proceedings)"
        ]
    },
    {
        "after_title": "Measurement Report-based Handover (AI-Enhanced Approaches)",
        "title": "Game‑Theoretic RL: Nash‑SAC",
        "bullets": [
            "Models multi‑user interactions as a game",
            "Computes Nash‑consistent policies with SAC",
            "Improved stability under heavy load/mobility",
            "Targets: drones/aircraft/ground terminals",
            "Source: Nash‑SAC LEO HO (arXiv:2402.00091)"
        ]
    },
    {
        "after_title": "Non-Measurement Report-based Handover (Traditional)",
        "title": "Location‑based HO with Particle Filter",
        "bullets": [
            "Fuses ephemeris + UE motion; robust to noise",
            "Particle filter tracks geometry for trigger timing",
            "Integrates with D2/CHO to reduce RLFs",
            "Hybrid with RL for adaptive thresholds",
            "Ref: Location‑based HO + PF (MDPI Electronics)"
        ]
    },
    {
        "after_title": "Non-Measurement Report-based Handover (Traditional)",
        "title": "Graph‑based Planning Details",
        "bullets": [
            "Time‑expanded graph: nodes=satellite/beam/time",
            "Edges weighted by link quality/HO cost",
            "Solvers: shortest path, min‑cost flow",
            "Complexity O(N^3 T); prune by visibility windows",
            "Refs: NTN optimization literature"
        ]
    },
    {
        "after_title": "Non-Measurement Report-based Handover (AI-Enhanced)",
        "title": "Transformer‑based Channel Prediction",
        "bullets": [
            "Attention captures long‑range temporal deps",
            "Inputs: RSRP history, Doppler, geometry",
            "Horizon: 10–30 s for proactive HO",
            "Feeds MPC or RL for look‑ahead",
            "Refs: LEO channel prediction (Transformer)"
        ]
    },
    {
        "after_title": "Non-Measurement Report-based Handover (AI-Enhanced)",
        "title": "CHO Trigger Prediction (XGBoost)",
        "bullets": [
            "Supervised model predicts success of CHO",
            "Features: geometry, timing, recent KPIs",
            "Reduces failed HOs and signaling",
            "Combines with D2 guardrails",
            "Ref: Two‑Step XGBoost for CHO (2023)"
        ]
    },
    {
        "after_title": "Performance Comparison Analysis",
        "title": "Evaluation Methodology",
        "bullets": [
            "KPIs: HO rate, RLF, throughput, E2E latency, QoS%",
            "Scenarios: urban, maritime, aerial (fast UEs)",
            "Constellations: Walker/mega‑LEO layouts",
            "Baselines: A4/A5/D2, CHO; vs RL/MPC/graphs",
            "Refs: NTN/LEO surveys 2023–2024"
        ]
    },
    {
        "after_title": "Performance Comparison Analysis",
        "title": "Implementation Guidelines",
        "bullets": [
            "Default to D2 + CHO; enable A4/A5 with tuned Hys/TTT",
            "Use offsets (Ofn/Ocn) to shape neighbor preference",
            "Predict visibility/Doppler for proactive HO",
            "Deploy RL incrementally with fallback/guardrails",
            "Cite: 3GPP TS 38.331; CHO (Rel‑16); RL papers"
        ]
    },
    {
        "after_title": "Introduction to LEO Satellite Networks",
        "title": "ISL/Backhaul Constraints & Mobility",
        "bullets": [
            "ISL changes alter routing, impact HO timing",
            "Backhaul bottlenecks can negate RSRP gains",
            "Joint HO + backhaul‑aware selection",
            "Plan handovers around ISL reconfiguration",
            "Refs: NTN architecture surveys (2024)"
        ]
    }
]


def _set_times_20(tf):
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(FONT_SIZE_PT)


def _find_slide_index_by_title(prs: Presentation, title_substr: str) -> int:
    for i, s in enumerate(prs.slides):
        t = s.shapes.title.text.strip() if s.shapes.title and s.shapes.title.has_text_frame else ""
        if title_substr.lower() in t.lower():
            return i
    return -1


def _insert_slide(prs: Presentation, insert_index: int, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    new_slide = prs.slides.add_slide(layout)
    new_slide.shapes.title.text = title
    ph = new_slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    # first bullet
    first = True
    for b in bullets:
        if first:
            tf.text = b
            first = False
        else:
            p = tf.add_paragraph()
            p.text = b
    _set_times_20(tf)
    # Reorder: move the new slide id to insert_index
    sldIdLst = prs.slides._sldIdLst  # lxml element list
    new_sldId = sldIdLst[-1]
    sldIdLst.remove(new_sldId)
    sldIdLst.insert(insert_index, new_sldId)


if __name__ == "__main__":
    base = Path("lession_expanded.pptx")
    if not base.exists():
        base = Path("lession.pptx")
    prs = Presentation(str(base))

    # Determine target insert positions by title, collect as (index_after, payload)
    planned = []
    for spec in SLIDES_TO_INSERT:
        idx = _find_slide_index_by_title(prs, spec["after_title"])  # index of the matched slide
        if idx >= 0:
            planned.append((idx + 1, spec))  # insert AFTER this index
    # Sort by insert index descending to avoid reindexing issues
    planned.sort(key=lambda x: x[0], reverse=True)

    for insert_after_idx, spec in planned:
        _insert_slide(prs, insert_after_idx, spec["title"], spec["bullets"])

    out = "lession_expanded2.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides")

