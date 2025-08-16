from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def add_slide_from_section(prs: Presentation, title: str, bullets: list[str], font_name: str = "Times New Roman", font_size_pt: int = 20):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # content placeholder is usually index 1
    placeholder = slide.placeholders[1]
    tf = placeholder.text_frame
    tf.clear()
    # First bullet initializes text_frame.text
    first_lvl, first_text = bullets[0]
    p = tf.paragraphs[0]
    p.level = first_lvl
    p.text = first_text
    for lvl, text in bullets[1:]:
        para = tf.add_paragraph()
        para.level = lvl
        para.text = text
    # Apply font to all paragraphs and runs
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)


def md_sections(md_text: str):
    sections = []
    current_title = None
    current_bullets: list[tuple[int, str]] = []
    for line in md_text.splitlines():
        if line.startswith('# '):
            if current_title and current_bullets:
                sections.append((current_title, current_bullets))
            current_title = line[2:].strip()
            current_bullets = []
        elif line.strip().startswith('- '):
            text = line.strip()[2:].strip()
            indent = (len(line) - len(line.lstrip(' '))) // 2
            current_bullets.append((indent, text))
        elif line.strip().startswith('* '):
            text = line.strip()[2:].strip()
            indent = (len(line) - len(line.lstrip(' '))) // 2
            current_bullets.append((indent, text))
        else:
            continue
    if current_title and current_bullets:
        sections.append((current_title, current_bullets))
    return sections

if __name__ == "__main__":
    src = "lession.pptx"
    md_path = "lession.md"
    out = "lession_expanded.pptx"
    prs = Presentation(src)
    md = Path(md_path).read_text(encoding='utf-8')
    sections = md_sections(md)
    # Only add the last section we appended (our new slide)
    if sections:
        title, bullets = sections[-1]
        add_slide_from_section(prs, title, bullets)
    prs.save(out)
    print(f"Wrote {out}")

