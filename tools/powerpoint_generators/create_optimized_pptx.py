#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os
import re

def set_mixed_font_style(text_frame, chinese_font="標楷體", english_font="Times New Roman", font_size=14):
    """設置混合中英文的字體樣式 - 正確處理每個字符"""
    for paragraph in text_frame.paragraphs:
        # 分析文本並分別設置字體
        text = paragraph.text
        if text:
            # 清除現有格式
            paragraph.clear()
            
            # 逐字符分析並設置字體
            i = 0
            while i < len(text):
                char = text[i]
                # 判斷是否為英文字符（包含數字和標點）
                if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', char):
                    # 英文字符，收集連續的英文
                    j = i
                    while j < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[j]):
                        j += 1
                    run = paragraph.add_run()
                    run.text = text[i:j]
                    run.font.name = english_font
                    run.font.size = Pt(font_size)
                    i = j
                else:
                    # 中文字符，收集連續的中文
                    j = i
                    while j < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[j]):
                        j += 1
                    run = paragraph.add_run()
                    run.text = text[i:j]
                    run.font.name = chinese_font
                    run.font.size = Pt(font_size)
                    i = j

def calculate_content_limits():
    """計算投影片內容區域的實際限制"""
    # PowerPoint 標準投影片尺寸 (16:9)
    slide_width = Inches(10)   # 10 英吋寬
    slide_height = Inches(7.5) # 7.5 英吋高
    
    # 內容區域（扣除標題和邊距）
    content_top = Inches(1.2)    # 標題下方
    content_bottom = Inches(7.2) # 底部邊距
    content_height = content_bottom - content_top  # 實際內容高度 = 6 英吋
    
    # 字型大小與行高計算
    font_size = 14  # 14pt 字型
    line_height = font_size * 1.2  # 行距為字型大小的 1.2 倍
    
    # 計算最大行數
    max_lines = int((content_height.pt) / line_height)
    
    print(f"投影片內容區域高度: {content_height}")
    print(f"字型大小: {font_size}pt")
    print(f"行高: {line_height}pt")
    print(f"最大行數: {max_lines}")
    
    return max_lines

def split_content_by_lines(content_text, max_lines_per_slide=22):
    """將長內容按行數分割成多個投影片"""
    lines = content_text.split('\n')
    slides_content = []
    current_slide = []
    current_line_count = 0
    
    for line in lines:
        # 計算這一行實際占用的行數（考慮長行自動換行）
        estimated_lines = max(1, len(line) // 80 + 1)  # 每80字符換行
        
        if current_line_count + estimated_lines > max_lines_per_slide:
            # 當前投影片已滿，開始新投影片
            if current_slide:
                slides_content.append('\n'.join(current_slide))
                current_slide = []
                current_line_count = 0
        
        current_slide.append(line)
        current_line_count += estimated_lines
    
    # 添加最後一個投影片
    if current_slide:
        slides_content.append('\n'.join(current_slide))
    
    return slides_content

def create_optimized_leo_presentation():
    """創建高度優化的 LEO 衛星 MC-HO 簡報"""
    
    # 計算內容限制
    max_lines = calculate_content_limits()
    print(f"每頁最大行數限制: {max_lines}")
    
    # 載入模板
    if os.path.exists('template.pptx'):
        prs = Presentation('template.pptx')
        print("使用模板 template.pptx")
    else:
        prs = Presentation()
        print("使用預設模板")
    
    # 清除現有投影片
    slide_count = len(prs.slides)
    for i in range(slide_count - 1, -1, -1):
        if i < len(prs.slides):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    
    # 定義版型
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # ====== 投影片 1: 標題頁 ======
    slide1 = prs.slides.add_slide(title_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "LEO 衛星網路多連線換手技術"
    subtitle1.text = """基於多連線與條件式換手方法提升換手性能
Multi-Connectivity and Conditional Handover Approach

論文來源：IEEE 2024
作者：Mohammed Al-Ansi 等
機構：University of Luxembourg SnT Centre"""
    
    set_mixed_font_style(title1.text_frame, font_size=22)
    set_mixed_font_style(subtitle1.text_frame, font_size=14)
    
    # ====== 投影片 2: 簡報大綱 ======
    slide2 = prs.slides.add_slide(content_layout)
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "簡報大綱 (Presentation Outline)"
    content2.text = """1. 研究背景與問題 (Background & Problem)
   • LEO 衛星網路挑戰
   • 傳統換手限制

2. MC-HO 演算法 (Core Algorithm)
   • 雙重連線架構
   • 條件式觸發機制

3. 演算法流程 (Algorithm Flow)
   • 4-Phase 設計
   • 時序圖分析

4. 數學模型 (Mathematical Model)
   • 距離計算公式
   • 路徑損耗模型

5. 實驗結果 (Experimental Results)
   • 性能比較分析

6. 實施與展望 (Implementation & Future)"""
    
    set_mixed_font_style(title2.text_frame, font_size=16)
    set_mixed_font_style(content2.text_frame, font_size=14)
    
    # ====== 投影片 3: LEO 衛星挑戰 ======
    slide3 = prs.slides.add_slide(content_layout)
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "LEO 衛星網路核心挑戰"
    content3.text = """• 系統參數特性
  - 軌道高度：600 km
  - 移動速度：7.56 km/s  
  - 波束直徑：50 km
  - 駐留時間：約 7 秒

• 傳統硬換手問題
  - Break-before-make 機制
  - 服務中斷：150ms
  - 頻繁換手：247 次/秒 (40% 重疊)

• 技術挑戰
  - 高速移動導致都卜勒頻移
  - 多衛星同頻干擾
  - 重疊區域連線不穩定"""
    
    set_mixed_font_style(title3.text_frame, font_size=16)
    set_mixed_font_style(content3.text_frame, font_size=14)
    
    # ====== 投影片 4: MC-HO 核心概念 ======
    slide4 = prs.slides.add_slide(content_layout)
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "MC-HO 核心概念"
    content4.text = """• 雙重連線架構
  - Master Node (MN)：服務衛星 SSAT
  - Secondary Node (SN)：目標衛星 TSAT
  - 同時維持兩條連線

• 觸發條件
  SC-HO: d_TSAT ≤ d_SSAT - d_offset
  MC-HO: d_TSAT ≤ R_b - d_offset AND 
         d_SSAT ≤ R_b - d_offset

• Make-before-break 流程
  1. 預建目標連線
  2. 確認穩定後釋放原連線
  3. 無服務中斷"""
    
    set_mixed_font_style(title4.text_frame, font_size=16)
    set_mixed_font_style(content4.text_frame, font_size=14)
    
    # ====== 投影片 5: 系統架構圖 ======
    slide5 = prs.slides.add_slide(content_layout)
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "MC-HO 系統架構"
    content5.text = """    ┌─────────────┐
    │ Core Network│
    └──────┬──────┘
           │
    ┌──────▼──────┐
    │     AMF     │
    └──────┬──────┘
           │
   ┌───────┼───────┐
   │       │       │
┌──▼──┐   │   ┌───▼─┐
│SSAT │◄─X2──►│TSAT │
│(MN) │       │(SN) │
└──┬──┘       └──┬──┘
   │             │
   │ ┌─────────┐ │
   └►│   UE    │◄┘
     │(Multi-  │
     │ Conn)   │
     └─────────┘"""
    
    set_mixed_font_style(title5.text_frame, font_size=16)
    set_mixed_font_style(content5.text_frame, font_size=12)
    
    # ====== 投影片 6: Phase 1-2 演算法流程 ======
    slide6 = prs.slides.add_slide(content_layout)
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "演算法流程 Phase 1-2"
    content6.text = """Phase 1: 初始化與監控
┌─────────────────────────────────┐
│ START → UE 連接 SSAT → 設為 MN  │
│   ↓                             │
│ 監控 UE 位置 (GNSS)             │
│   ↓                             │
│ 計算與候選衛星距離               │
│   ↓                             │
│ 持續評估 CHO 條件               │
└─────────────────────────────────┘

Phase 2: 條件式換手觸發
┌─────────────────────────────────┐
│ IF (d_TSAT ≤ R_b - d_offset)    │
│ AND (d_SSAT ≤ R_b - d_offset)   │
│   ↓                             │
│ CHO 條件滿足                    │
│   ↓                             │
│ 觸發 SN 新增程序                │
│   ↓                             │
│ 發送測量報告至 SSAT             │
└─────────────────────────────────┘"""
    
    set_mixed_font_style(title6.text_frame, font_size=16)
    set_mixed_font_style(content6.text_frame, font_size=12)
    
    # ====== 投影片 7: Phase 3-4 演算法流程 ======
    slide7 = prs.slides.add_slide(content_layout)
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "演算法流程 Phase 3-4"
    content7.text = """Phase 3: 雙重連線建立
┌─────────────────────────────────┐
│ SSAT 發送 SN Addition Request   │
│   ↓                             │
│ TSAT 接受並回應                 │
│   ↓                             │
│ UE 執行 Random Access 至 TSAT   │
│   ↓                             │
│ 建立 UE-TSAT 連線               │
│   ↓                             │
│ 啟動 Packet Duplication         │
└─────────────────────────────────┘

Phase 4: 路徑切換與釋放
┌─────────────────────────────────┐
│ TSAT 傳送首個封包               │
│   ↓                             │
│ 發送 Path Switch Request        │
│   ↓                             │
│ AMF 執行 Bearer Modification    │
│   ↓                             │
│ UPF 重新路由流量                │
│   ↓                             │
│ 釋放 SSAT 連線                  │
└─────────────────────────────────┘"""
    
    set_mixed_font_style(title7.text_frame, font_size=16)
    set_mixed_font_style(content7.text_frame, font_size=12)
    
    # ====== 投影片 8: 封包複製機制 ======
    slide8 = prs.slides.add_slide(content_layout)
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "封包複製機制 (Packet Duplication)"
    content8.text = """• 觸發條件
  CHO 條件滿足 → 啟動封包複製
  Master Cell Group + Split Bearer

• 資料流路徑
  Path 1: Core Network → SSAT → UE
  Path 2: Core Network → SSAT → TSAT → UE

• 選擇合併 (Selection Combining)
  - 即時監控兩路徑 SINR
  - 選擇最高 SINR 路徑
  - SINR_sel = max(SINR_SSAT, SINR_TSAT)

• 分集增益
  - 典型增益：2-5 dB
  - 降低邊緣用戶連線失效
  - 提升系統可靠性"""
    
    set_mixed_font_style(title8.text_frame, font_size=16)
    set_mixed_font_style(content8.text_frame, font_size=14)
    
    # ====== 投影片 9: 時序圖 ======
    slide9 = prs.slides.add_slide(content_layout)
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "MC-HO 時序圖 (Message Sequence)"
    content9.text = """UE     SSAT    TSAT    AMF     UPF
│       │       │       │       │
│◄─測量─►│       │       │       │
│ 報告   │       │       │       │
│       │◄─SN──►│       │       │
│       │ Add   │       │       │
│◄─RRC──┤       │       │       │
│Reconf │       │       │       │
│◄─RA───────────►│       │       │
│       │       │       │       │
│◄══雙重資料流═══►│       │       │
│     (PD)      │       │       │
│       │       │◄Path─►│       │
│       │       │Switch │       │
│       │       │       │◄Mod──►│
│       │       │       │Bearer │
│◄─UE Context───│       │       │
│  Release      │       │       │
│◄═══單一資料流═══►│       │       │"""
    
    set_mixed_font_style(title9.text_frame, font_size=16)
    set_mixed_font_style(content9.text_frame, font_size=11)
    
    # ====== 投影片 10: 核心演算法偽代碼 (第1部分) ======
    slide10 = prs.slides.add_slide(content_layout)
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "核心演算法偽代碼 (Part 1)"
    content10.text = """Algorithm: MC-HO_Decision_Engine
Input: UE_pos, SSAT_pos, TSAT_pos
Output: Handover_decision

1  BEGIN
2    // 初始化 Initialization
3    d_SSAT ← calc_distance(UE, SSAT)
4    d_TSAT ← calc_distance(UE, TSAT)
5    R_beam ← 25000  // meters
6    d_offset ← 1000 // meters
7    
8    // CHO 條件檢查
9    IF (d_TSAT ≤ R_beam - d_offset) AND
10      (d_SSAT ≤ R_beam - d_offset) THEN
11   
12      // SN 新增程序
13      send_SN_request(SSAT → TSAT)
14      establish_UE_TSAT_connection()
15      activate_packet_duplication()"""
    
    set_mixed_font_style(title10.text_frame, font_size=16)
    set_mixed_font_style(content10.text_frame, font_size=12)
    
    # ====== 投影片 11: 核心演算法偽代碼 (第2部分) ======
    slide11 = prs.slides.add_slide(content_layout)
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "核心演算法偽代碼 (Part 2)"
    content11.text = """      // 選擇合併迴圈
16      WHILE (in_overlap_region) DO
17        SINR_SSAT ← measure_SINR(SSAT)
18        SINR_TSAT ← measure_SINR(TSAT)
19        threshold ← 3  // dB
20        
21        IF (SINR_TSAT > SINR_SSAT 
22                      + threshold) THEN
23          select_primary_path(TSAT)
24        ELSE
25          select_primary_path(SSAT)
26        END IF
27      END WHILE
28      
29      // 路徑切換
30      execute_path_switch(TSAT_as_MN)
31      release_connection(SSAT)
32   END IF
33 END"""
    
    set_mixed_font_style(title11.text_frame, font_size=16)
    set_mixed_font_style(content11.text_frame, font_size=12)
    
    # ====== 投影片 12: 數學模型與公式 ======
    slide12 = prs.slides.add_slide(content_layout)
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "數學模型與公式 (Mathematical Model)"
    content12.text = """• 距離計算 (論文公式 4)
  d = √(R_E² sin²(α) + h₀² + 2h₀·R_E) - R_E sin(α)
  R_E = 6371 km, h₀ = 600 km

• 接收功率 (論文公式 1)
  R_UE = EIRP - PL_total (dBm)

• 總路徑損耗 (論文公式 2)
  PL_total = Pr_LOS × PL_LOS + 
             (1 - Pr_LOS) × PL_NLOS

• LOS 路徑損耗 (論文公式 3)
  PL_LOS = 32.45 + 20log₁₀(f_c) + 
           20log₁₀(d) + SF (dB)

• 觸發條件比較
  SC-HO: d_TSAT ≤ d_SSAT - d_offset
  MC-HO: d_TSAT ≤ R_b - d_offset AND
         d_SSAT ≤ R_b - d_offset"""
    
    set_mixed_font_style(title12.text_frame, font_size=16)
    set_mixed_font_style(content12.text_frame, font_size=12)
    
    # ====== 投影片 13: 實驗結果 - 換手次數 ======
    slide13 = prs.slides.add_slide(content_layout)
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "實驗結果：換手次數比較"
    content13.text = """• 換手性能 vs 波束重疊
┌─────┬─────┬─────┬─────┐
│重疊%│SC-HO│MC-HO│改善%│
├─────┼─────┼─────┼─────┤
│  0  │ 148 │ 148 │  0  │
│ 10  │ 165 │ 162 │ 1.8 │
│ 20  │ 185 │ 145 │21.6 │
│ 30  │ 212 │ 129 │39.2 │
│ 40  │ 247 │ 130 │47.4 │
└─────┴─────┴─────┴─────┘

• 關鍵發現
  - 40% 重疊時效果最佳
  - 換手次數減少近一半
  - MC-HO 延遲換手執行
  - 避免 ping-pong 效應

• 技術原理
  MC-HO 在多覆蓋區域維持雙重連線，
  直到接近波束邊緣才執行換手。"""
    
    set_mixed_font_style(title13.text_frame, font_size=16)
    set_mixed_font_style(content13.text_frame, font_size=13)
    
    # ====== 投影片 14: 實驗結果 - RLF 分析 ======
    slide14 = prs.slides.add_slide(content_layout)
    title14 = slide14.shapes.title
    content14 = slide14.placeholders[1]
    
    title14.text = "實驗結果：連結失效分析"
    content14.text = """• RLF 性能比較
┌─────┬─────┬─────┬─────┐
│重疊%│SC-HO│MC-HO│減少%│
├─────┼─────┼─────┼─────┤
│  0  │ 168 │ 168 │  0  │
│ 10  │ 221 │ 211 │ 4.5 │
│ 20  │ 296 │ 265 │10.5 │
│ 30  │ 403 │ 338 │16.1 │
│ 40  │ 532 │ 410 │22.9 │
└─────┴─────┴─────┴─────┘

• RLF 減少機制
  1. 傳輸分集：雙路徑冗餘保護
  2. 路徑選擇：選擇最佳 SINR
  3. 邊緣保護：避免信號衰弱
  4. 干擾緩解：選擇合併效果

• 分集增益：2-5 dB 典型值"""
    
    set_mixed_font_style(title14.text_frame, font_size=16)
    set_mixed_font_style(content14.text_frame, font_size=13)
    
    # ====== 投影片 15: 系統容量分析 ======
    slide15 = prs.slides.add_slide(content_layout)
    title15 = slide15.shapes.title
    content15 = slide15.placeholders[1]
    
    title15.text = "系統容量分析"
    content15.text = """• 平均容量 vs 重疊 (Mb/s/Hz)
┌─────┬─────┬─────┬─────┐
│重疊%│SC-HO│MC-HO│提升%│
├─────┼─────┼─────┼─────┤
│  0  │ 2.1 │ 2.1 │  0  │
│ 10  │ 2.0 │ 2.05│ 2.5 │
│ 20  │ 1.9 │ 2.0 │ 5.3 │
│ 30  │ 1.8 │ 1.95│ 8.3 │
│ 40  │ 1.7 │ 1.9 │11.8 │
└─────┴─────┴─────┴─────┘

• 容量提升原理
  - 傳輸分集增益提升 SINR
  - 選擇合併減少干擾
  - 負載分散避免過載

• 重要觀察
  重疊增加雖造成干擾，但 MC-HO
  仍保持容量優勢。分集增益有效
  補償干擾損失。"""
    
    set_mixed_font_style(title15.text_frame, font_size=16)
    set_mixed_font_style(content15.text_frame, font_size=13)
    
    # ====== 投影片 16: 複雜度分析 ======
    slide16 = prs.slides.add_slide(content_layout)
    title16 = slide16.shapes.title
    content16 = slide16.placeholders[1]
    
    title16.text = "複雜度與成本分析"
    content16.text = """• 計算複雜度
┌──────────┬─────┬─────┬─────┐
│ 項目       │SC-HO│MC-HO│倍數 │
├──────────┼─────┼─────┼─────┤
│ 距離計算   │O(1) │O(2) │ 2x  │
│ 連線管理   │O(1) │O(2) │ 2x  │
│ 信令處理   │基準 │+25% │1.25x│
│ 封包處理   │基準 │+50% │1.5x │
└──────────┴─────┴─────┴─────┘

• 資源消耗
  UE: 記憶體 +30%, 功耗 +15%
  衛星: 處理負載 +25%
  地面網路: AMF/UPF 負載 +15-20%

• 效益成本比
  成本增加: 20-30%
  性能改善: 換手-47%, RLF-23%
  ROI: 1.5-2 年回收"""
    
    set_mixed_font_style(title16.text_frame, font_size=16)
    set_mixed_font_style(content16.text_frame, font_size=12)
    
    # ====== 投影片 17: 實施挑戰 ======
    slide17 = prs.slides.add_slide(content_layout)
    title17 = slide17.shapes.title
    content17 = slide17.placeholders[1]
    
    title17.text = "實施挑戰與解決方案"
    content17.text = """• 主要挑戰
  1. 時間同步：±1μs 精度要求
     解決：GPS 統一時基 + ISL 輔助

  2. 干擾管理：同頻干擾增加 3-5dB
     解決：功率控制 + 波束成型

  3. 負載平衡：雙連線負載不均
     解決：智慧調度 + 動態分配

• 標準化進展
  3GPP Release 17: NTN 基礎 (2022)
  3GPP Release 18: MC-HO 研究 (2024)
  3GPP Release 19: 預期標準化 (2025)

• 部署考量
  - 衛星星座規模：100+ 衛星
  - 用戶終端：需多連線能力
  - 初期投資 vs 長期效益"""
    
    set_mixed_font_style(title17.text_frame, font_size=16)
    set_mixed_font_style(content17.text_frame, font_size=12)
    
    # ====== 投影片 18: 未來研究方向 ======
    slide18 = prs.slides.add_slide(content_layout)
    title18 = slide18.shapes.title
    content18 = slide18.placeholders[1]
    
    title18.text = "未來研究方向"
    content18.text = """• 演算法增強
  - AI/ML 增強決策
  - 機器學習預測用戶移動
  - 多準則決策引擎

• 系統擴展
  - 大規模星座 (1000+ 衛星)
  - 分層決策架構
  - 異質網路融合 (LEO/MEO/GEO)

• 6G 整合
  - 網路切片專用換手策略
  - 邊緣計算輔助決策
  - 數位孿生網路最佳化

• 應用場景
  - 高速移動載具
  - 海事航空通訊
  - 災害應急通訊
  - IoT 與 M2M 通訊"""
    
    set_mixed_font_style(title18.text_frame, font_size=16)
    set_mixed_font_style(content18.text_frame, font_size=13)
    
    # ====== 投影片 19: 結論 ======
    slide19 = prs.slides.add_slide(content_layout)
    title19 = slide19.shapes.title
    content19 = slide19.placeholders[1]
    
    title19.text = "結論 (Conclusions)"
    content19.text = """• 主要貢獻
  ✓ 首次提出位置基準 MC-HO 演算法
  ✓ 結合條件式換手與多連線技術
  ✓ 4-Phase 架構清晰可實作

• 顯著性能提升
  ✓ 換手次數減少最高 47.4%
  ✓ RLF 減少最高 22.9%
  ✓ 系統容量提升最高 11.8%

• 實用價值
  - 為 LEO 營運商提供實施方案
  - 推動 5G/6G NTN 標準化
  - 改善衛星通訊服務品質

• 關鍵訊息
  MC-HO 透過智慧多連線管理和位置
  感知觸發，有效解決 LEO 衛星網路
  換手問題，為 6G 衛星系統奠定基礎。"""
    
    set_mixed_font_style(title19.text_frame, font_size=16)
    set_mixed_font_style(content19.text_frame, font_size=13)
    
    # ====== 投影片 20: Q&A ======
    slide20 = prs.slides.add_slide(content_layout)
    title20 = slide20.shapes.title
    content20 = slide20.placeholders[1]
    
    title20.text = "Q&A 與討論 (Questions & Discussion)"
    content20.text = """• 常見技術問題
  Q1: UE 硬體需求為何？
  A1: 雙重 RRC、多天線、增強基頻處理

  Q2: 三顆以上衛星多連線？
  A2: 1MN + 多 SN 架構，分層決策

  Q3: 高負載下性能？
  A3: 智慧調度維持性能，需准入控制

• 實施問題
  Q4: 5G 核心網整合？
  A4: 主要 AMF/UPF 軟體升級

  Q5: 部署成本？
  A5: 初期 +20-30%，1.5-2年回收

  Q6: 頻譜效率影響？
  A6: 整體提升 8-12%

謝謝聆聽！歡迎討論"""
    
    set_mixed_font_style(title20.text_frame, font_size=16)
    set_mixed_font_style(content20.text_frame, font_size=13)
    
    # 儲存簡報
    output_filename = "LEO衛星MC-HO演算法最佳化簡報.pptx"
    prs.save(output_filename)
    print(f"最佳化簡報已成功創建：{output_filename}")
    print(f"總投影片數：{len(prs.slides)}")
    
    return output_filename

if __name__ == "__main__":
    create_optimized_leo_presentation()