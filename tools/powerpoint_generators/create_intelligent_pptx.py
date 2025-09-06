#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智慧型 PowerPoint 生成器 - 整合圖表分析與教學內容
Intelligent PowerPoint Generator with Figure Analysis Integration
"""

import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from intelligent_figure_selector import IntelligentFigureSelector

class IntelligentPowerPointGenerator:
    """智慧型 PowerPoint 生成器，整合圖表分析功能"""
    
    def __init__(self, template_path="../../doc/template.pptx"):
        self.template_path = template_path
        self.max_lines_per_slide = 20
        self.figure_selector = IntelligentFigureSelector()
        
        # 圖片資源路徑
        self.image_base_path = "../../論文圖片"  # 論文圖片實際位置
        
        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"
        
    def load_template(self):
        """載入簡報模板"""
        try:
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                print(f"✅ 模板載入成功: {self.template_path}")
            else:
                prs = Presentation()
                print("⚠️  模板不存在，使用預設模板")
            return prs
        except Exception as e:
            print(f"❌ 模板載入失敗: {e}")
            return Presentation()

    def set_mixed_font_style(self, text_frame, font_size=14):
        """設定中英文混合字體"""
        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()
                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', char):
                        # 英文字符 - 收集連續的英文
                        eng_text = ""
                        while i < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            eng_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = eng_text
                        run.font.name = self.english_font
                        run.font.size = Pt(font_size)
                    else:
                        # 中文字符 - 收集連續的中文
                        chn_text = ""
                        while i < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            chn_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = chn_text
                        run.font.name = self.chinese_font
                        run.font.size = Pt(font_size)

    def estimate_content_lines(self, content_list):
        """估算內容行數"""
        total_lines = 0
        for item in content_list:
            lines = item.count('\n') + 1
            # 每 60 個字符約占一行
            char_lines = len(item) // 60 + (1 if len(item) % 60 > 0 else 0)
            total_lines += max(lines, char_lines)
        return total_lines

    def split_content_for_slides(self, content_list, max_lines=20):
        """將內容分割為適合投影片的片段"""
        slides_content = []
        current_slide_content = []
        current_lines = 0
        
        for item in content_list:
            item_lines = self.estimate_content_lines([item])
            
            if current_lines + item_lines <= max_lines:
                current_slide_content.append(item)
                current_lines += item_lines
            else:
                if current_slide_content:
                    slides_content.append(current_slide_content)
                current_slide_content = [item]
                current_lines = item_lines
        
        if current_slide_content:
            slides_content.append(current_slide_content)
        
        return slides_content

    def add_image_to_slide(self, slide, image_path, left=None, top=None, width=None, height=None):
        """將圖片添加到投影片"""
        try:
            if left is None:
                left = Inches(6.5)
            if top is None:
                top = Inches(1.5)
            if width is None:
                width = Inches(3)
            if height is None:
                height = Inches(2.5)
            
            if os.path.exists(image_path):
                picture = slide.shapes.add_picture(image_path, left, top, width, height)
                print(f"✅ 圖片已添加: {os.path.basename(image_path)}")
                return picture
            else:
                print(f"⚠️  圖片不存在: {image_path}")
                return None
        except Exception as e:
            print(f"❌ 圖片添加失敗: {e}")
            return None

    def create_title_slide(self, prs):
        """創建標題投影片"""
        slide_layout = prs.slide_layouts[0]  # 標題版面
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "LEO 衛星網路多連線換手技術"
        self.set_mixed_font_style(title.text_frame, font_size=24)
        
        subtitle_content = """基於 Multi-Connectivity 和 Conditional Handover 的性能提升方法

🛰️ Multi-Connectivity Handover (MC-HO)
📡 3GPP NTN 標準應用
🔄 Soft Handover 技術實現
📊 性能驗證與分析

本簡報基於 IEEE 2024 論文
「Enhancing Handover Performance in LEO Satellite Networks」"""
        
        subtitle.text = subtitle_content
        self.set_mixed_font_style(subtitle.text_frame, font_size=16)

    def create_figure_slides(self, prs, selected_figures, explanations):
        """根據選中的重要圖表創建投影片"""
        
        for figure_info in selected_figures:
            figure_name = figure_info['figure_name']
            description = figure_info['description']
            image_file = figure_info['image_file']
            
            # 獲取詳細說明內容
            explanation = explanations.get(figure_name, "")
            
            # 分割內容以適應投影片高度
            content_lines = explanation.split('\n')
            slides_content = self.split_content_for_slides(content_lines, self.max_lines_per_slide)
            
            for i, slide_content in enumerate(slides_content):
                slide_layout = prs.slide_layouts[1]  # 標題與內容版面
                slide = prs.slides.add_slide(slide_layout)
                
                # 設定標題
                if len(slides_content) > 1:
                    title_text = f"{figure_name}: {description} ({i+1}/{len(slides_content)})"
                else:
                    title_text = f"{figure_name}: {description}"
                
                slide.shapes.title.text = title_text
                self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
                
                # 設定內容
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_text = '\n'.join(slide_content)
                    content_placeholder.text = content_text
                    self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)
                
                # 在第一張投影片添加圖片 (如果存在)
                if i == 0 and image_file:
                    image_path = os.path.join(self.image_base_path, image_file)
                    if not os.path.exists(image_path):
                        # 嘗試在其他路徑查找圖片
                        alt_paths = [
                            f"../../論文圖片/{image_file}",
                            f"論文圖片/{image_file}",
                            image_file
                        ]
                        for alt_path in alt_paths:
                            if os.path.exists(alt_path):
                                image_path = alt_path
                                break
                    
                    self.add_image_to_slide(slide, image_path)

    def create_conclusion_slide(self, prs):
        """創建結論投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "結論與主要貢獻"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
        
        conclusion_content = """🎯 主要技術貢獻：

📡 MC-HO 演算法創新：
• 整合 Multi-Connectivity 與 Conditional Handover
• 實現真正的 Soft Handover 機制
• Location-based 觸發條件優化

⚡ 顯著性能改善：
• 換手次數減少達 47% (40% overlap 條件下)
• 無線連結失效降低 23%
• 保持更佳的系統容量

🛰️ LEO 衛星網路適用性：
• 適應高速移動特性 (7.56 km/s)
• 充分利用多重覆蓋區域
• 提升服務連續性與可靠性

🔄 未來發展方向：
• Quasi-Earth Fixed Beams 應用
• 更高頻段 (Ka-band) 擴展  
• Multi-criteria 觸發條件研究
• 智慧排程算法優化"""
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = conclusion_content
            self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

    def generate_intelligent_presentation(self, output_filename="LEO衛星MC-HO演算法智慧簡報.pptx"):
        """生成智慧型簡報"""
        
        print("🚀 開始生成智慧型 PowerPoint 簡報...")
        print("="*60)
        
        # 1. 載入模板
        prs = self.load_template()
        
        # 2. 智慧選擇重要圖表
        print("🎯 執行智慧圖表選擇...")
        selected_figures = self.figure_selector.select_figures_by_priority(max_figures=5)
        explanations = self.figure_selector.generate_figure_explanations(selected_figures)
        
        print(f"✅ 已選擇 {len(selected_figures)} 個重要圖表")
        for fig in selected_figures:
            print(f"   • {fig['figure_name']}: {fig['description']}")
        
        # 3. 創建投影片內容
        print("\n📄 開始創建投影片...")
        
        # 標題投影片
        self.create_title_slide(prs)
        print("✅ 標題投影片已創建")
        
        # 圖表解釋投影片
        self.create_figure_slides(prs, selected_figures, explanations)
        print(f"✅ {len(selected_figures)} 個圖表投影片已創建")
        
        # 結論投影片
        self.create_conclusion_slide(prs)
        print("✅ 結論投影片已創建")
        
        # 4. 儲存簡報
        output_path = f"../../doc/{output_filename}"
        try:
            prs.save(output_path)
            print(f"\n🎉 簡報生成成功！")
            print(f"📁 輸出檔案: {output_path}")
            print(f"📊 總投影片數: {len(prs.slides)}")
            
            # 生成統計報告
            self.generate_creation_report(selected_figures, len(prs.slides))
            
        except Exception as e:
            print(f"❌ 簡報儲存失敗: {e}")

    def generate_creation_report(self, selected_figures, total_slides):
        """生成簡報創建報告"""
        
        report = f"""# 📊 智慧型 PowerPoint 簡報創建報告

## 🎯 簡報概覽
- **總投影片數**: {total_slides} 張
- **核心技術圖表**: {len(selected_figures)} 個
- **創建時間**: 2024-09-06
- **模式**: 智慧圖表選擇 + 詳細技術解釋

## 📋 圖表選擇結果
"""
        
        for i, fig in enumerate(selected_figures, 1):
            priority_stars = "⭐" * fig['teaching_value']
            report += f"### {i}. {fig['figure_name']}\n"
            report += f"- **教學價值**: {priority_stars}\n"
            report += f"- **說明**: {fig['description']}\n"
            report += f"- **圖片**: {fig['image_file'] or '需手動添加'}\n\n"

        report += f"""## ✅ 技術特色
- **智慧圖表選擇**: 基於教學價值與技術重要性自動篩選
- **詳細技術解釋**: 每個圖表配有深度技術說明
- **中英文字體混合**: 標楷體 + Times New Roman 
- **投影片高度控制**: 最多 {self.max_lines_per_slide} 行，自動分頁
- **圖文對應**: 重要圖表與詳細說明精確對應

## 🎓 教學完整性
- ✅ 系統模型建立 (Figure 1)
- ✅ 核心演算法詳解 (Figure 2) 
- ✅ 實驗參數說明 (Table 1)
- ✅ 性能驗證分析 (Figure 3, 4)
- ✅ 結論與貢獻總結

## 📈 相比之前版本的改進
- **有針對性**: 只選擇最重要的圖表，避免資訊過載
- **有深度**: 每個圖表都有詳細的技術解釋
- **有邏輯**: 從基礎概念到性能驗證的完整教學流程
- **有品質**: 圖表與說明內容精確對應，提升理解效果
"""

        with open("../../doc/智慧簡報創建報告.md", 'w', encoding='utf-8') as f:
            f.write(report)
        
        print("📝 創建報告已儲存至: ../doc/智慧簡報創建報告.md")

def main():
    """主程式"""
    generator = IntelligentPowerPointGenerator()
    generator.generate_intelligent_presentation()

if __name__ == "__main__":
    main()