#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
3GPP TS 38.331 換手事件機制投影片生成器
根據ts.md內容製作A4、A5、D2事件的技術說明投影片
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

class ThreeGPPHandoverEventsPresentation:
    """3GPP換手事件機制投影片生成器"""

    def __init__(self, output_dir="../../doc"):
        self.output_dir = output_dir

        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"

        # 顏色設定
        self.primary_color = RGBColor(0, 51, 102)      # 深藍色
        self.secondary_color = RGBColor(255, 102, 0)   # 橙色
        self.formula_color = RGBColor(139, 0, 139)     # 深紫色

        # 確保輸出目錄存在
        os.makedirs(self.output_dir, exist_ok=True)

    def set_mixed_font_style(self, text_frame, chinese_font=None, english_font=None, font_size=12):
        """設置混合中英文字體"""
        if chinese_font is None:
            chinese_font = self.chinese_font
        if english_font is None:
            english_font = self.english_font

        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()

                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', char):
                        # 收集連續的英文字符
                        j = i
                        while j < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = english_font
                        run.font.size = Pt(font_size)
                        i = j
                    else:
                        # 收集連續的中文字符
                        j = i
                        while j < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = chinese_font
                        run.font.size = Pt(font_size)
                        i = j

    def create_3gpp_handover_events_slide(self):
        """創建3GPP換手事件機制投影片"""
        print("🚀 開始創建3GPP TS 38.331換手事件機制投影片...")

        # 創建簡報
        prs = Presentation()

        # 創建單一投影片展示3個換手機制
        slide_layout = prs.slide_layouts[1]  # 標題與內容版面
        slide = prs.slides.add_slide(slide_layout)

        # 設定標題
        slide.shapes.title.text = "3GPP TS 38.331 標準換手事件機制"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)

        # 創建內容
        content_text = """【Event A4：鄰居小區優於門檻值】
觸發條件：Mn + Ofn + Ocn – Hys > Thresh
離開條件：Mn + Ofn + Ocn + Hys < Thresh

• Mn：鄰居小區測量結果 (RSRP: dBm, RSRQ/RS-SINR: dB)
• Ofn：鄰居小區測量對象特定偏移值 (dB)
• Ocn：鄰居小區特定偏移值 (dB，未配置時為0)
• Hys：遲滞參數 (dB)
• Thresh：事件門檻參數 (與Mn相同單位)

【Event A5：服務小區劣於門檻1且鄰居小區優於門檻2】
觸發條件：Mp + Hys < Thresh1 AND Mn + Ofn + Ocn – Hys > Thresh2
離開條件：Mp – Hys > Thresh1 OR Mn + Ofn + Ocn + Hys < Thresh2

• Mp：NR SpCell測量結果，不考慮任何偏移
• 其他變數定義同Event A4
• Thresh1：服務小區門檻 (與Mp相同單位)
• Thresh2：鄰居小區門檻 (與Mn相同單位)

【Event D2：UE與服務小區移動參考位置距離超過門檻1且與移動參考位置距離低於門檻2】
觸發條件：Ml1 – Hys > Thresh1 AND Ml2 + Hys < Thresh2
離開條件：Ml1 + Hys < Thresh1 OR Ml2 – Hys > Thresh2

• Ml1：UE與服務小區移動參考位置的距離 (米)
• Ml2：UE與目標移動參考位置的距離 (米)
• 移動參考位置基於衛星星曆和epoch時間確定
• 所有距離參數以米為單位"""

        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content_text
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=11)

        # 儲存簡報
        output_filename = os.path.join(self.output_dir, "3GPP_TS38331_換手事件機制.pptx")
        prs.save(output_filename)

        print(f"✅ 3GPP換手事件機制投影片創建完成")
        print(f"📊 投影片數：1張")
        print(f"📁 儲存位置：{output_filename}")

        # 生成創建報告
        self.generate_creation_report()

        return output_filename

    def generate_creation_report(self):
        """生成創建報告"""
        report_content = """# 3GPP TS 38.331 換手事件機制投影片創建報告

## 📊 投影片概覽
- **主題**: 3GPP TS 38.331 標準換手事件機制
- **投影片數**: 1張
- **創建時間**: 2024-09-17
- **內容來源**: doc/ts.md

## 🎯 包含的換手事件

### Event A4：鄰居小區優於門檻值
- **用途**: 當鄰居小區信號品質超過設定門檻時觸發
- **核心公式**: Mn + Ofn + Ocn – Hys > Thresh
- **適用場景**: 主動發現更好的服務小區

### Event A5：雙門檻判斷機制
- **用途**: 服務小區劣化且鄰居小區優良時觸發
- **核心公式**:
  - 服務小區：Mp + Hys < Thresh1
  - 鄰居小區：Mn + Ofn + Ocn – Hys > Thresh2
- **適用場景**: 避免不必要的換手，確保換手品質

### Event D2：基於距離的NTN換手
- **用途**: 專為衛星通訊設計的位置感知換手
- **核心公式**:
  - 離開條件：Ml1 – Hys > Thresh1
  - 接近條件：Ml2 + Hys < Thresh2
- **適用場景**: LEO衛星網路的地理位置感知換手

## 🔬 技術特點

### 變數定義完整性
- 所有公式變數都有明確定義
- 單位規範（dBm, dB, 米）
- 參數來源清楚標示

### 條件邏輯清晰
- 觸發條件與離開條件分別說明
- 邏輯運算符（AND, OR）明確標示
- 遲滯機制防止乒乓效應

### NTN特殊考量
- Event D2專為衛星通訊設計
- 考慮移動參考位置和衛星星曆
- 距離測量取代傳統信號強度

## ✅ 投影片品質
- **技術準確性**: 直接引用3GPP TS 38.331標準
- **內容完整性**: 涵蓋觸發/離開條件和所有變數
- **排版清晰**: 結構化展示，便於理解
- **實用價值**: 可直接用於技術培訓或標準解釋

## 📈 應用價值
此投影片適合用於：
- 3GPP標準培訓
- NTN技術講解
- 換手演算法對比
- 系統設計參考

---
*本投影片基於3GPP TS 38.331 version 18.5.1 Release 18標準文件製作*"""

        report_filename = os.path.join(self.output_dir, "3GPP換手事件機制投影片報告.md")
        with open(report_filename, 'w', encoding='utf-8') as f:
            f.write(report_content)

        print(f"📋 創建報告已生成：{report_filename}")

def main():
    """主函數"""
    print("🚀 3GPP TS 38.331 換手事件機制投影片生成器")
    print("=" * 50)

    # 創建生成器實例
    generator = ThreeGPPHandoverEventsPresentation()

    # 生成投影片
    try:
        output_file = generator.create_3gpp_handover_events_slide()
        print(f"\n✅ 投影片生成成功：{output_file}")
        print("\n🎯 投影片特色：")
        print("   • 完整涵蓋3個重要換手事件")
        print("   • 精確的公式表達和變數定義")
        print("   • 結構化的內容排版")
        print("   • 基於官方3GPP標準")

    except Exception as e:
        print(f"❌ 投影片生成失敗：{e}")
        return False

    return True

if __name__ == "__main__":
    main()