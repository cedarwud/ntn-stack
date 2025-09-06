#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import re
import json

def set_mixed_font_style(text_frame, chinese_font="標楷體", english_font="Times New Roman", font_size=14):
    """設置混合中英文字體"""
    for paragraph in text_frame.paragraphs:
        text = paragraph.text
        if text:
            paragraph.clear()
            
            i = 0
            while i < len(text):
                char = text[i]
                if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', char):
                    j = i
                    while j < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[j]):
                        j += 1
                    run = paragraph.add_run()
                    run.text = text[i:j]
                    run.font.name = english_font
                    run.font.size = Pt(font_size)
                    i = j
                else:
                    j = i
                    while j < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[j]):
                        j += 1
                    run = paragraph.add_run()
                    run.text = text[i:j]
                    run.font.name = chinese_font
                    run.font.size = Pt(font_size)
                    i = j

def add_image_to_slide(slide, image_path, left=None, top=None, width=None, height=None):
    """在投影片中添加圖片"""
    try:
        if not os.path.exists(image_path):
            print(f"⚠️  圖片不存在：{image_path}")
            return None
            
        # 預設位置和大小
        if left is None:
            left = Inches(1)
        if top is None:
            top = Inches(2)
        if width is None:
            width = Inches(8)
        if height is None:
            height = Inches(5)
            
        # 添加圖片
        picture = slide.shapes.add_picture(image_path, left, top, width, height)
        print(f"✅ 已添加圖片：{os.path.basename(image_path)}")
        return picture
        
    except Exception as e:
        print(f"❌ 添加圖片失敗：{e}")
        return None

def create_image_slide(prs, title, image_path, description="", layout_idx=8):
    """創建包含圖片的投影片"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = title
        set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
        
        # 添加圖片 (調整位置以適應標題)
        if os.path.exists(image_path):
            add_image_to_slide(
                slide, 
                image_path, 
                left=Inches(0.5), 
                top=Inches(1.5), 
                width=Inches(9), 
                height=Inches(5.5)
            )
        
        # 如果有描述，添加到底部
        if description:
            # 創建文字框
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
            textbox.text = description
            set_mixed_font_style(textbox.text_frame, font_size=12)
        
        return slide
        
    except Exception as e:
        print(f"❌ 創建圖片投影片失敗：{e}")
        return None

def load_extraction_info():
    """載入圖片提取資訊"""
    try:
        with open("pdf_extraction_summary.json", 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"⚠️  無法載入提取資訊：{e}")
        return None

def create_enhanced_presentation_with_images():
    """創建包含原始論文圖片的增強版簡報"""
    
    print("🚀 開始創建包含原始論文圖片的增強版簡報")
    
    # 載入模板
    if os.path.exists('template.pptx'):
        prs = Presentation('template.pptx')
        print("✅ 使用模板 template.pptx")
    else:
        prs = Presentation()
        print("⚠️  使用預設模板")
    
    # 清除現有投影片
    slide_count = len(prs.slides)
    for i in range(slide_count - 1, -1, -1):
        if i < len(prs.slides):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    
    # 載入提取的圖片資訊
    extraction_info = load_extraction_info()
    
    # 定義版型
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    picture_layout = prs.slide_layouts[8] if len(prs.slide_layouts) > 8 else content_layout
    
    print("📝 開始創建投影片...")
    
    # ====== 投影片 1: 標題頁 ======
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "LEO 衛星網路多連線換手技術"
    slide1.placeholders[1].text = """基於多連線與條件式換手方法提升換手性能
Multi-Connectivity and Conditional Handover

論文：IEEE 2024 • 作者：Mohammed Al-Ansi 等
機構：University of Luxembourg SnT Centre"""
    
    set_mixed_font_style(slide1.shapes.title.text_frame, font_size=22)
    set_mixed_font_style(slide1.placeholders[1].text_frame, font_size=16)
    
    # ====== 投影片 2: LEO 衛星系統模型 (論文 Figure 1) ======
    figure1_path = "論文圖片/page_2_img_1.png"
    if os.path.exists(figure1_path):
        slide2 = create_image_slide(
            prs,
            "LEO 衛星系統模型 (論文 Figure 1)",
            figure1_path,
            "圖片來源：論文 Figure 1 - LEO satellites coverage scenario"
        )
    else:
        # 備用方案：使用文字描述
        slide2 = prs.slides.add_slide(content_layout)
        slide2.shapes.title.text = "LEO 衛星系統模型"
        slide2.placeholders[1].text = """• 系統配置參數
  - 軌道高度：600 km
  - 衛星速度：7.56 km/s
  - 波束直徑：50 km
  - 多覆蓋區域設計

• 技術挑戰
  - 高速移動導致頻繁換手
  - 波束駐留時間短 (~7 秒)
  - 多衛星同頻干擾"""
        set_mixed_font_style(slide2.shapes.title.text_frame, font_size=18)
        set_mixed_font_style(slide2.placeholders[1].text_frame, font_size=15)
    
    # ====== 投影片 3: MC-HO 演算法核心概念 ======
    slide3 = prs.slides.add_slide(content_layout)
    slide3.shapes.title.text = "MC-HO 核心概念與架構"
    slide3.content = """• 雙重連線架構 (Dual Connectivity)
  - Master Node (MN)：服務衛星 SSAT
  - Secondary Node (SN)：目標衛星 TSAT

• 觸發條件改進
  傳統 SC-HO: d_TSAT ≤ d_SSAT - d_offset
  提出 MC-HO: d_TSAT ≤ R_b - d_offset AND
              d_SSAT ≤ R_b - d_offset

• Make-before-break 機制
  1. 預先建立目標連線
  2. 封包複製與選擇合併
  3. 無縫路徑切換"""
    slide3.placeholders[1].text = slide3.content
    set_mixed_font_style(slide3.shapes.title.text_frame, font_size=18)
    set_mixed_font_style(slide3.placeholders[1].text_frame, font_size=14)
    
    # ====== 投影片 4: MC-HO 程序流程圖 (論文 Figure 2) ======
    # 使用第4頁的流程圖 - 可能是 Figure 2
    slide4 = prs.slides.add_slide(content_layout)
    slide4.shapes.title.text = "MC-HO 演算法流程"
    slide4.content = """Phase 1: 初始化與監控
• UE 連接 SSAT (設為 MN)
• GNSS 位置監控與距離計算

Phase 2: 條件式觸發
• CHO 條件評估
• 多覆蓋區域檢測

Phase 3: 雙重連線建立  
• SN Addition Request/Response
• Random Access Procedure
• Packet Duplication 啟動

Phase 4: 路徑切換
• Path Switch Request
• Bearer Modification
• 連線釋放與角色轉換"""
    slide4.placeholders[1].text = slide4.content
    set_mixed_font_style(slide4.shapes.title.text_frame, font_size=18)
    set_mixed_font_style(slide4.placeholders[1].text_frame, font_size=13)
    
    # ====== 投影片 5: 實驗結果 - 換手次數比較 (論文 Figure 3) ======
    figure3_path = "論文圖片/page_4_img_1.png"  # 可能是 Figure 3
    if os.path.exists(figure3_path):
        slide5 = create_image_slide(
            prs,
            "實驗結果：換手次數比較 (論文 Figure 3)",
            figure3_path,
            "MC-HO vs SC-HO 換手性能比較 - 40% 重疊時改善 47.4%"
        )
    else:
        slide5 = prs.slides.add_slide(content_layout)
        slide5.shapes.title.text = "實驗結果：換手次數比較"
        slide5.content = """• 不同波束重疊比例的性能比較
┌─────┬─────┬─────┬─────┐
│重疊%│SC-HO│MC-HO│改善%│
├─────┼─────┼─────┼─────┤
│  0  │ 148 │ 148 │  0  │
│ 10  │ 165 │ 162 │ 1.8 │
│ 20  │ 185 │ 145 │21.6 │
│ 30  │ 212 │ 129 │39.2 │
│ 40  │ 247 │ 130 │47.4 │
└─────┴─────┴─────┴─────┘

• 關鍵發現
  - 重疊比例越高，MC-HO 優勢越顯著
  - 40% 重疊時換手次數減少近一半"""
        slide5.placeholders[1].text = slide5.content
        set_mixed_font_style(slide5.shapes.title.text_frame, font_size=18)
        set_mixed_font_style(slide5.placeholders[1].text_frame, font_size=13)
    
    # ====== 投影片 6: 無線連結失效分析 (論文 Figure 4) ======
    figure4_path = "論文圖片/page_4_img_2.png"  # 可能是 Figure 4
    if os.path.exists(figure4_path):
        slide6 = create_image_slide(
            prs,
            "無線連結失效分析 (論文 Figure 4)",
            figure4_path,
            "RLF 性能比較 - MC-HO 最高減少 22.9% 連結失效"
        )
    
    # ====== 投影片 7: 系統容量分析 (論文 Figure 6) ======
    figure6_path = "論文圖片/page_5_img_2.png"  # 可能是 Figure 6
    if os.path.exists(figure6_path):
        slide7 = create_image_slide(
            prs,
            "系統容量分析 (論文 Figure 6)",
            figure6_path,
            "平均容量 vs 重疊百分比 - MC-HO 保持較高容量"
        )
    
    # ====== 投影片 8: 時間序列分析 (論文 Figure 5) ======
    figure5_path = "論文圖片/page_5_img_1.png"  # 可能是 Figure 5
    if os.path.exists(figure5_path):
        slide8 = create_image_slide(
            prs,
            "時間序列分析 (論文 Figure 5)",
            figure5_path,
            "40% 重疊場景下的換手行為比較 - MC-HO 更平穩一致"
        )
    
    # ====== 投影片 9: 數學模型與公式 ======
    slide9 = prs.slides.add_slide(content_layout)
    slide9.shapes.title.text = "數學模型與公式推導"
    slide9.content = """• 距離計算公式 (論文公式 4)
  d = √(R_E² sin²(α) + h₀² + 2h₀·R_E) - R_E sin(α)

• 接收功率 (論文公式 1)
  R_UE = EIRP - PL_total (dBm)

• 路徑損耗 (論文公式 2-3)
  PL_total = Pr_LOS × PL_LOS + (1-Pr_LOS) × PL_NLOS
  PL_LOS = 32.45 + 20log₁₀(f_c) + 20log₁₀(d) + SF

• MC-HO 觸發條件 (論文公式 7)
  d_TSAT(t) ≤ R_b - d_offset AND
  d_SSAT(t) ≤ R_b - d_offset

參數：R_E = 6371km, h₀ = 600km, R_b = 25km"""
    slide9.placeholders[1].text = slide9.content
    set_mixed_font_style(slide9.shapes.title.text_frame, font_size=18)
    set_mixed_font_style(slide9.placeholders[1].text_frame, font_size=12)
    
    # ====== 投影片 10: 結論與貢獻 ======
    slide10 = prs.slides.add_slide(content_layout)
    slide10.shapes.title.text = "結論與主要貢獻"
    slide10.content = """• 技術創新
  ✓ 首次提出位置基準 MC-HO 演算法
  ✓ 整合條件式換手與多連線技術
  ✓ 4-Phase 架構設計清晰可實作

• 顯著性能提升
  ✓ 換手次數減少：最高 47.4% (40% 重疊)
  ✓ 連結失效減少：最高 22.9%
  ✓ 系統容量提升：8-12%

• 實用價值
  - 為 LEO 星座營運商提供實施方案
  - 推動 5G/6G NTN 標準化發展  
  - 改善全球衛星通訊服務品質

• 未來展望
  - AI/ML 增強的智慧決策
  - 大規模星座部署最佳化
  - 6G 異質網路融合"""
    slide10.placeholders[1].text = slide10.content
    set_mixed_font_style(slide10.shapes.title.text_frame, font_size=18)
    set_mixed_font_style(slide10.placeholders[1].text_frame, font_size=13)
    
    # 儲存簡報
    output_filename = "LEO衛星MC-HO演算法簡報-含原始圖表.pptx"
    prs.save(output_filename)
    print(f"\n✅ 包含原始論文圖表的簡報已創建完成！")
    print(f"📄 檔案名稱：{output_filename}")
    print(f"📊 總投影片數：{len(prs.slides)}")
    
    # 顯示使用的圖片清單
    print(f"\n📷 使用的原始論文圖片：")
    used_images = [
        ("Figure 1", "論文圖片/page_2_img_1.png", "LEO 衛星覆蓋場景"),
        ("Figure 3", "論文圖片/page_4_img_1.png", "換手次數比較"),
        ("Figure 4", "論文圖片/page_4_img_2.png", "RLF 性能分析"),
        ("Figure 5", "論文圖片/page_5_img_1.png", "時間序列分析"),
        ("Figure 6", "論文圖片/page_5_img_2.png", "系統容量分析")
    ]
    
    for fig_name, path, description in used_images:
        if os.path.exists(path):
            print(f"   ✅ {fig_name}: {description}")
        else:
            print(f"   ⚠️  {fig_name}: 圖片未找到")
    
    return output_filename

if __name__ == "__main__":
    create_enhanced_presentation_with_images()