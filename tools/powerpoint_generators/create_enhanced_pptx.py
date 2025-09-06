#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import os

def set_font_style(text_frame, chinese_font="標楷體", english_font="Times New Roman", font_size=16):
    """設置文字框的字體樣式"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = chinese_font
        paragraph.font.size = Pt(font_size)
        # 為英文設置不同字體
        for run in paragraph.runs:
            run.font.name = english_font if any(c.isalpha() and ord(c) < 128 for c in run.text) else chinese_font
            run.font.size = Pt(font_size)

def create_enhanced_leo_presentation():
    """創建增強版 LEO 衛星網路多連線換手技術教學簡報"""
    
    # 載入模板
    if os.path.exists('template.pptx'):
        prs = Presentation('template.pptx')
        print("使用模板 template.pptx")
    else:
        prs = Presentation()
        print("使用預設模板")
    
    # 清除模板中的現有投影片
    slide_count = len(prs.slides)
    for i in range(slide_count - 1, -1, -1):
        if i < len(prs.slides):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    
    # 定義版型
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    comparison_layout = prs.slide_layouts[4]
    
    # ====== 投影片 1: 標題頁 ======
    slide1 = prs.slides.add_slide(title_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "LEO 衛星網路多連線換手技術"
    subtitle1.text = """基於多連線與條件式換手方法提升換手性能
(Multi-Connectivity and Conditional Handover Approach)

論文來源：IEEE 2024
作者：Mohammed Al-Ansi 等
機構：University of Luxembourg SnT Centre"""
    
    set_font_style(title1.text_frame, font_size=24)
    set_font_style(subtitle1.text_frame, font_size=16)
    
    # ====== 投影片 2: 核心問題與挑戰 ======
    slide2 = prs.slides.add_slide(content_layout)
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "LEO 衛星網路的核心挑戰"
    content2.text = """• LEO 衛星特性
  - 軌道高度：600 km
  - 移動速度：7.56 km/s
  - 波束直徑：50 km
  - 波束駐留時間：僅約 7 秒

• 傳統硬換手 (Hard Handover) 的問題
  - Break-before-make 機制
  - 服務中斷時間：平均 150ms
  - 無線連結失效率高
  - 頻繁換手：每秒高達 247 次 (40% 重疊時)

• 關鍵技術挑戰
  - 高速移動導致都卜勒頻移
  - 多衛星同頻干擾
  - 用戶在重疊區域的連線不穩定"""
    
    set_font_style(title2.text_frame, font_size=20)
    set_font_style(content2.text_frame, font_size=16)
    
    # ====== 投影片 3: MC-HO 演算法架構 ======
    slide3 = prs.slides.add_slide(content_layout)
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "MC-HO 演算法架構 (Algorithm Architecture)"
    content3.text = """• 雙重連線架構 (Dual Connectivity)
  - Master Node (MN)：服務衛星 SSAT
  - Secondary Node (SN)：目標衛星 TSAT
  - 同時維持兩條連線路径

• 條件式換手觸發機制 (Conditional Handover)
  - 位置基準觸發條件：
    d_TSAT(t) ≤ R_b - d_offset  AND  d_SSAT(t) ≤ R_b - d_offset
  - R_b = 25 km (波束半徑)
  - d_offset = 1~5 km (距離偏移)

• Make-before-break 軟換手流程
  1. 預先建立目標連線
  2. 確認連線穩定後釋放原連線
  3. 無服務中斷時間"""
    
    set_font_style(title3.text_frame, font_size=20)
    set_font_style(content3.text_frame, font_size=16)
    
    # ====== 投影片 4: 核心演算法流程圖 ======
    slide4 = prs.slides.add_slide(content_layout)
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "MC-HO 核心演算法流程"
    content4.text = """Phase 1: 初始化與監控
├── UE 連接至 SSAT (設為 MN)
├── 持續監控用戶位置 (GNSS)
└── 計算與候選衛星距離

Phase 2: 觸發條件評估
├── 檢查多覆蓋區域條件
├── IF (d_TSAT ≤ R_b - d_offset) AND (d_SSAT ≤ R_b - d_offset)
└── THEN 觸發 SN 新增程序

Phase 3: SN 新增與封包複製
├── SSAT 發送 SN Addition Request 至 TSAT
├── UE 執行 Random Access 至 TSAT
├── 建立 UE-TSAT 連線
├── 啟動 Packet Duplication (PD)
└── 雙路徑同步傳輸開始

Phase 4: 路徑切換與釋放
├── TSAT 傳送首個資料封包
├── 發送 MN Change Request 至 AMF
├── AMF 執行 Bearer Modification
├── UPF 重新路由核心網路流量
└── 釋放 SSAT 連線，TSAT 成為新 MN"""
    
    set_font_style(title4.text_frame, font_size=18)
    set_font_style(content4.text_frame, font_size=14)
    
    # ====== 投影片 5: 封包複製機制詳解 ======
    slide5 = prs.slides.add_slide(content_layout)
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "封包複製機制 (Packet Duplication) 詳解"
    content5.text = """• 封包複製啟動條件
  - CHO 條件滿足時立即啟動
  - Master Cell Group with Split Bearer 模式
  - 核心網路連接至 MN，流量分別傳至 MN 和 SN

• 資料流路徑
  Path 1: Core Network → SSAT (MN) → UE
  Path 2: Core Network → SSAT → TSAT (SN) → UE

• 選擇比合併 (Selection Combining)
  - 即時監控兩路徑 SINR 值
  - 動態選擇最高 SINR 路徑
  - SINR_selected = max(SINR_SSAT, SINR_TSAT)

• 分集增益計算
  - 典型分集增益：2-5 dB
  - 有效降低邊緣用戶連線失效
  - 提升整體系統可靠性"""
    
    set_font_style(title5.text_frame, font_size=18)
    set_font_style(content5.text_frame, font_size=16)
    
    # ====== 投影片 6: 數學模型與公式推導 ======
    slide6 = prs.slides.add_slide(content_layout)
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "數學模型與公式推導"
    content6.text = """• 距離計算公式
  d = √(R_E² sin²(α) + h₀² + 2h₀ · R_E) - R_E sin(α)
  其中：R_E = 6371 km, h₀ = 600 km, α = 仰角

• 路徑損耗模型 (3GPP NTN Standard)
  R_UE = EIRP - PL_total
  PL_total = Pr_LOS × PL_LOS + (1 - Pr_LOS) × PL_NLOS

• LOS 路徑損耗
  PL_LOS = 32.45 + 20log₁₀(f_c) + 20log₁₀(d) + SF

• NLOS 路徑損耗
  PL_NLOS = 32.45 + 20log₁₀(f_c) + 20log₁₀(d) + SF + CL

• MC-HO vs SC-HO 觸發條件比較
  SC-HO: d_TSAT(t) ≤ d_SSAT(t) - d_offset
  MC-HO: d_TSAT(t) ≤ R_b - d_offset AND d_SSAT(t) ≤ R_b - d_offset"""
    
    set_font_style(title6.text_frame, font_size=18)
    set_font_style(content6.text_frame, font_size=14)
    
    # ====== 投影片 7: 實驗設計與參數配置 ======
    slide7 = prs.slides.add_slide(content_layout)
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "實驗設計與參數配置"
    content7.text = """• 仿真環境參數
┌─────────────────────┬─────────────┐
│ 參數項目                    │ 數值        │
├─────────────────────┼─────────────┤
│ LEO 衛星高度               │ 600 km      │
│ 衛星速度                   │ 7.56 km/s   │
│ 衛星 Tx 最大增益           │ 30 dBi      │
│ 波束直徑                   │ 50 km       │
│ EIRP 密度                  │ 34 dBW/MHz  │
│ 載波頻率                   │ 2 GHz       │
│ 頻寬                       │ 30 MHz      │
│ 用戶密度                   │ 1 user/km²  │
│ 距離偏移                   │ 1 km / 5 km │
│ 仿真時間                   │ 200 秒      │
└─────────────────────┴─────────────┘

• 評估指標
  - 平均換手次數/秒 (Average Handover Rate)
  - 無線連結失效次數/秒 (Radio Link Failure Rate)
  - 平均系統容量 (Average System Capacity)"""
    
    set_font_style(title7.text_frame, font_size=18)
    set_font_style(content7.text_frame, font_size=14)
    
    # ====== 投影片 8: 性能比較結果分析 (1/2) ======
    slide8 = prs.slides.add_slide(content_layout)
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "性能比較結果分析 (1/2) - 換手次數"
    content8.text = """• 不同波束重疊比例的換手性能比較
┌────────┬─────────┬─────────┬─────────┐
│重疊比例    │SC-HO(次/秒)│MC-HO(次/秒)│改善率    │
├────────┼─────────┼─────────┼─────────┤
│   0%       │    148      │    148      │    0%    │
│  10%       │    165      │    162      │   1.8%   │
│  20%       │    185      │    145      │  21.6%   │
│  30%       │    212      │    129      │  39.2%   │
│  40%       │    247      │    130      │  47.4%   │
└────────┴─────────┴─────────┴─────────┘

• 關鍵發現
  - 重疊比例越高，MC-HO 優勢越明顯
  - 40% 重疊時，換手次數減少近一半
  - MC-HO 在多覆蓋區域延遲換手執行
  - 有效避免 "ping-pong" 效應"""
    
    set_font_style(title8.text_frame, font_size=18)
    set_font_style(content8.text_frame, font_size=14)
    
    # ====== 投影片 9: 性能比較結果分析 (2/2) ======
    slide9 = prs.slides.add_slide(content_layout)
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "性能比較結果分析 (2/2) - 連結失效"
    content9.text = """• 無線連結失效 (RLF) 比較
┌────────┬─────────┬─────────┬─────────┐
│重疊比例    │SC-HO(次/秒)│MC-HO(次/秒)│改善率    │
├────────┼─────────┼─────────┼─────────┤
│   0%       │    168      │    168      │    0%    │
│  10%       │    221      │    211      │   4.5%   │
│  20%       │    296      │    265      │  10.5%   │
│  30%       │    403      │    338      │  16.1%   │
│  40%       │    532      │    410      │  22.9%   │
└────────┴─────────┴─────────┴─────────┘

• RLF 減少的技術原因
  1. 傳輸分集 (Transmit Diversity)
     - 雙路徑同時傳輸提供冗餘保護
  2. 最佳連結選擇 (Best Link Selection)  
     - 即時選擇最高 SINR 路徑
  3. 邊緣用戶保護 (Cell Edge Protection)
     - 避免單一衛星信號衰弱影響"""
    
    set_font_style(title9.text_frame, font_size=18)
    set_font_style(content9.text_frame, font_size=14)
    
    # ====== 投影片 10: 演算法複雜度分析 ======
    slide10 = prs.slides.add_slide(content_layout)
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "演算法複雜度分析"
    content10.text = """• 計算複雜度 (Computational Complexity)
  - SC-HO: O(1) - 簡單距離比較
  - MC-HO: O(2) - 雙重連線管理

• 信令開銷 (Signaling Overhead)
┌─────────────┬──────┬──────┬─────────┐
│ 信令類型        │ SC-HO  │ MC-HO  │ 增加幅度  │
├─────────────┼──────┼──────┼─────────┤
│ 測量報告        │ 基準   │ +10%   │ 輕微      │
│ 連線建立        │ 基準   │ +100%  │ 顯著      │
│ 資料傳輸        │ 基準   │ +50%   │ 中等      │
│ 連線釋放        │ 基準   │ +20%   │ 輕微      │
└─────────────┴──────┴──────┴─────────┘

• 記憶體與處理需求
  - UE 端：需支援雙重協議堆疊
  - 衛星端：處理能力需求提升 15-20%
  - 地面站：信令處理負載增加 25%

• 能耗分析
  - UE 功耗增加：約 20%
  - 系統總能耗：增加 10% (考慮換手減少效益)"""
    
    set_font_style(title10.text_frame, font_size=18)
    set_font_style(content10.text_frame, font_size=14)
    
    # ====== 投影片 11: 時序圖與訊息流程 ======
    slide11 = prs.slides.add_slide(content_layout)
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "MC-HO 時序圖與訊息流程"
    content11.text = """UE ←→ SSAT ←→ TSAT ←→ AMF ←→ UPF

1. 測量與報告階段
   UE → SSAT: Measurement Report
   SSAT: 評估 CHO 條件

2. SN 新增階段  
   SSAT → TSAT: SN Addition Request
   TSAT → SSAT: SN Addition Response
   UE ← SSAT: RRC Reconfiguration
   UE → TSAT: Random Access Procedure

3. 封包複製階段
   SSAT ↔ TSAT: Data Forwarding Setup
   UE ← SSAT & TSAT: 雙重資料流傳輸
   UE: Selection Combining (選擇最佳路徑)

4. 路徑切換階段
   TSAT → AMF: Path Switch Request  
   AMF → UPF: Modify Bearer Request
   UPF → AMF: Modify Bearer Response
   AMF → TSAT: Path Switch Acknowledge

5. 連線釋放階段
   TSAT → SSAT: UE Context Release Request
   SSAT: 釋放資源並清除 UE 上下文"""
    
    set_font_style(title11.text_frame, font_size=18)
    set_font_style(content11.text_frame, font_size=13)
    
    # ====== 投影片 12: 關鍵演算法偽代碼 ======
    slide12 = prs.slides.add_slide(content_layout)
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "關鍵演算法偽代碼 (Pseudocode)"
    content12.text = """Algorithm: MC-HO_Decision_Engine
Input: UE_position, SSAT_position, TSAT_position
Output: Handover_decision

BEGIN
  // 計算距離
  d_SSAT = calculate_distance(UE_position, SSAT_position)
  d_TSAT = calculate_distance(UE_position, TSAT_position)
  
  // MC-HO 觸發條件檢查
  IF (d_TSAT <= R_beam - d_offset) AND 
     (d_SSAT <= R_beam - d_offset) THEN
     
     // 啟動 SN 新增程序
     send_SN_addition_request(SSAT, TSAT)
     establish_UE_TSAT_connection()
     
     // 啟動封包複製
     activate_packet_duplication()
     
     WHILE (in_overlap_region) DO
        SINR_SSAT = measure_SINR(SSAT)
        SINR_TSAT = measure_SINR(TSAT)
        
        // 選擇最佳路徑
        IF (SINR_TSAT > SINR_SSAT + threshold) THEN
           select_path(TSAT)
        ELSE
           select_path(SSAT)
        END IF
     END WHILE
     
     // 執行路徑切換
     execute_path_switch(TSAT_as_new_MN)
     release_connection(SSAT)
  END IF
END"""
    
    set_font_style(title12.text_frame, font_size=18)
    set_font_style(content12.text_frame, font_size=12)
    
    # ====== 投影片 13: 系統架構與介面 ======
    slide13 = prs.slides.add_slide(content_layout)
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "MC-HO 系統架構與介面"
    content13.text = """• 系統組件架構
  ┌─────────────┐    ┌─────────────┐
  │     UE      │◄──►│    SSAT     │
  │ (Multi-conn)│    │    (MN)     │
  └─────────────┘    └─────────────┘
         │                   │
         │            ┌─────────────┐
         └───────────►│    TSAT     │
                      │    (SN)     │
                      └─────────────┘

• 關鍵介面與協定
  - Uu Interface: UE ↔ 衛星 (NR)
  - X2 Interface: SSAT ↔ TSAT (資料轉發)
  - S1-U Interface: 衛星 ↔ UPF (用戶平面)  
  - S1-C Interface: 衛星 ↔ AMF (控制平面)

• 協議堆疊修改
  - UE: 支援雙重 RRC 連線
  - SSAT/TSAT: X2AP 增強 (SN Addition/Release)
  - AMF: 路徑切換處理邏輯
  - UPF: 封包複製與路由切換

• 與 3GPP 標準對齊
  - 基於 Release 15/16 DC 架構
  - NTN-specific enhancements
  - 向後相容性保證"""
    
    set_font_style(title13.text_frame, font_size=18)
    set_font_style(content13.text_frame, font_size=14)
    
    # ====== 投影片 14: 實施挑戰與解決方案 ======
    slide14 = prs.slides.add_slide(content_layout)
    title14 = slide14.shapes.title
    content14 = slide14.placeholders[1]
    
    title14.text = "實施挑戰與解決方案"
    content14.text = """• 主要技術挑戰
  1. 同步問題 (Synchronization)
     挑戰：多衛星間時間同步精度要求
     解決：GPS/GNSS 統一時基 + ISL 輔助同步

  2. 干擾管理 (Interference Management)  
     挑戰：同頻多衛星造成干擾增加
     解決：功率控制 + 波束成型 + 頻率規劃

  3. 負載平衡 (Load Balancing)
     挑戰：雙重連線造成網路負載不均
     解決：智慧調度演算法 + 動態資源分配

• 部署考量因素
  - 衛星星座規模：適用於 100+ 衛星系統
  - 用戶終端改造：需支援多連線能力
  - 地面基礎設施：需升級核心網設備
  - 運營成本：初期投資 vs 長期效益

• 標準化進展
  - 3GPP Release 17: NTN 基礎標準
  - 3GPP Release 18: MC-HO 考慮中
  - ITU-R 研究項目：衛星間換手最佳化"""
    
    set_font_style(title14.text_frame, font_size=18)
    set_font_style(content14.text_frame, font_size=14)
    
    # ====== 投影片 15: 結論與貢獻 ======
    slide15 = prs.slides.add_slide(content_layout)
    title15 = slide15.shapes.title
    content15 = slide15.placeholders[1]
    
    title15.text = "結論與研究貢獻"
    content15.text = """• 主要技術貢獻
  1. 首次提出位置基準 MC-HO 演算法
     - 利用 LEO 軌道可預測特性
     - 條件式換手減少失效機率

  2. 建立完整數學模型與分析
     - 3GPP NTN 標準相容路徑損耗模型
     - 量化分析性能提升機制

  3. 實驗驗證顯著性能改善
     - 換手次數減少：最高 47.4%
     - 連結失效減少：最高 22.9%
     - 系統容量提升：8% 增益

• 實用價值與影響
  - 為 LEO 星座提供實用換手解決方案
  - 推動 5G/6G NTN 標準化進程  
  - 改善全球衛星通訊服務品質

• 未來研究方向
  - 大規模星座部署最佳化
  - AI/ML 增強的智慧換手決策
  - 異質衛星網路融合架構"""
    
    set_font_style(title15.text_frame, font_size=18)
    set_font_style(content15.text_frame, font_size=16)
    
    # ====== 投影片 16: Q&A ======
    slide16 = prs.slides.add_slide(content_layout)
    title16 = slide16.shapes.title
    content16 = slide16.placeholders[1]
    
    title16.text = "Questions & Discussion"
    content16.text = """• 技術問題討論
  Q1: MC-HO 對用戶終端硬體有何要求？
  Q2: 如何處理三顆以上衛星的多連線場景？
  Q3: 在高負載情況下的性能表現如何？

• 實施相關問題  
  Q4: 部署成本與傳統方案的比較？
  Q5: 與現有 5G 核心網的整合難度？
  Q6: 對頻譜效率的實際影響為何？

• 研究延伸方向
  Q7: 是否適用於 GEO/MEO 衛星？
  Q8: 與地面 5G 網路的協同可能性？
  Q9: 在移動用戶場景下的適用性？

感謝聆聽！
歡迎提問與討論"""
    
    set_font_style(title16.text_frame, font_size=18)
    set_font_style(content16.text_frame, font_size=16)
    
    # 儲存簡報
    output_filename = "LEO衛星網路MC-HO演算法詳解簡報.pptx"
    prs.save(output_filename)
    print(f"增強版簡報已成功創建：{output_filename}")
    print(f"總投影片數：{len(prs.slides)}")
    
    return output_filename

if __name__ == "__main__":
    create_enhanced_leo_presentation()