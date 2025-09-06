#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
以演算法為核心的 PowerPoint 生成器
Algorithm-Focused PowerPoint Generator for MC-HO LEO Satellite Networks
"""

import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

class AlgorithmFocusedPowerPointGenerator:
    """以演算法為核心的 PowerPoint 生成器"""
    
    def __init__(self, template_path="../../doc/template.pptx"):
        self.template_path = template_path
        self.max_lines_per_slide = 20
        
        # 圖片資源路徑 (輔助用)
        self.image_base_path = "../../論文圖片"
        
        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"
        
        # 演算法核心內容結構
        self.algorithm_structure = {
            "introduction": "LEO 衛星網路切換挑戰與 MC-HO 解決方案",
            "problem_definition": "Single Connectivity Handover 的問題分析",
            "algorithm_core": "Multi-Connectivity Handover 演算法核心設計",
            "technical_details": "MC-HO 技術實現細節",
            "performance_analysis": "演算法性能分析與驗證",
            "practical_implementation": "實際應用與部署考量"
        }
    
    def load_template(self):
        """載入簡報模板"""
        try:
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                print(f"✅ 模板載入成功: {self.template_path}")
            else:
                prs = Presentation()
                print("⚠️  模板不存在，使用預設模板")
            return prs
        except Exception as e:
            print(f"❌ 模板載入失敗: {e}")
            return Presentation()

    def estimate_content_lines(self, content_text):
        """估算內容所需行數 - 遵循技術指南"""
        lines = content_text.split('\n')
        total_lines = 0
        
        for line in lines:
            if not line.strip():  # 空行
                total_lines += 1
            else:
                # 長行自動換行計算 (每行最多 80 字符)
                char_count = len(line)
                estimated_lines = max(1, (char_count + 79) // 80)
                total_lines += estimated_lines
        
        return total_lines

    def split_content_for_slides(self, content_text, max_lines=20):
        """將內容分割為適合投影片的片段 - 遵循技術指南"""
        if self.estimate_content_lines(content_text) <= max_lines:
            return [content_text]
        
        lines = content_text.split('\n')
        parts = []
        current_part = []
        current_lines = 0
        
        for line in lines:
            line_count = max(1, (len(line) + 79) // 80) if line.strip() else 1
            
            if current_lines + line_count > max_lines and current_part:
                # 當前頁面已滿，開始新頁面
                parts.append('\n'.join(current_part))
                current_part = [line]
                current_lines = line_count
            else:
                current_part.append(line)
                current_lines += line_count
        
        # 添加最後一個頁面
        if current_part:
            parts.append('\n'.join(current_part))
        
        return parts

    def set_mixed_font_style(self, text_frame, font_size=14):
        """設定中英文混合字體"""
        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()
                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', char):
                        # 英文字符 - 收集連續的英文
                        eng_text = ""
                        while i < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            eng_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = eng_text
                        run.font.name = self.english_font
                        run.font.size = Pt(font_size)
                    else:
                        # 中文字符 - 收集連續的中文
                        chn_text = ""
                        while i < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            chn_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = chn_text
                        run.font.name = self.chinese_font
                        run.font.size = Pt(font_size)

    def add_supporting_image(self, slide, image_file, caption=""):
        """添加輔助說明圖片 (右上角，較小尺寸)"""
        if not image_file:
            return None
            
        image_path = self.find_image_path(image_file)
        if image_path:
            try:
                # 輔助圖片位置：右上角，較小尺寸
                left = Inches(7.5)
                top = Inches(1.2)
                width = Inches(2.2)
                height = Inches(1.5)
                
                picture = slide.shapes.add_picture(image_path, left, top, width, height)
                
                # 添加圖片說明
                if caption:
                    textbox = slide.shapes.add_textbox(
                        Inches(7.5), Inches(2.8), Inches(2.2), Inches(0.3)
                    )
                    text_frame = textbox.text_frame
                    text_frame.text = caption
                    self.set_mixed_font_style(text_frame, font_size=10)
                
                print(f"✅ 輔助圖片已添加: {os.path.basename(image_path)}")
                return picture
            except Exception as e:
                print(f"❌ 圖片添加失敗: {e}")
        return None

    def find_image_path(self, image_file):
        """尋找圖片檔案路徑"""
        possible_paths = [
            os.path.join(self.image_base_path, image_file),
            f"../../論文圖片/{image_file}",
            image_file
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None

    def create_algorithm_comparison_table(self, slide):
        """創建演算法比較表格"""
        # 添加比較表格
        table = slide.shapes.add_table(
            rows=6, cols=3, 
            left=Inches(1), top=Inches(2), 
            width=Inches(8), height=Inches(4)
        ).table
        
        # 表格標題行
        table.cell(0, 0).text = "比較項目"
        table.cell(0, 1).text = "SC-HO (傳統方法)"
        table.cell(0, 2).text = "MC-HO (提出方法)"
        
        # 表格內容
        comparison_data = [
            ("連線架構", "單一連線 (Single)", "雙重連線 (Dual)"),
            ("切換方式", "硬切換 (Break-before-make)", "軟切換 (Make-before-break)"),
            ("資料傳輸", "單一路徑", "封包複製 (Packet Duplication)"),
            ("觸發機制", "測量基礎 (SINR)", "位置基礎 + 條件式 (CHO)"),
            ("可靠性", "易中斷", "高可靠性 (Selection Combining)")
        ]
        
        for i, (item, sc_ho, mc_ho) in enumerate(comparison_data, 1):
            table.cell(i, 0).text = item
            table.cell(i, 1).text = sc_ho  
            table.cell(i, 2).text = mc_ho
        
        # 設定表格字體
        for row in range(6):
            for col in range(3):
                cell = table.cell(row, col)
                self.set_mixed_font_style(cell.text_frame, font_size=12)
                # 標題行使用粗體
                if row == 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

    def create_performance_metrics_table(self, slide):
        """創建性能指標表格"""
        table = slide.shapes.add_table(
            rows=5, cols=4,
            left=Inches(1), top=Inches(2.5),
            width=Inches(8), height=Inches(3.5)
        ).table
        
        # 表格標題
        headers = ["重疊百分比", "SC-HO 換手次數/秒", "MC-HO 換手次數/秒", "改善幅度"]
        for col, header in enumerate(headers):
            table.cell(0, col).text = header
        
        # 性能數據
        performance_data = [
            ("0%", "148", "148", "0%"),
            ("20%", "185", "145", "21.6%"),
            ("30%", "212", "129", "39.2%"),
            ("40%", "247", "130", "47.4%")
        ]
        
        for row, (overlap, sc, mc, improvement) in enumerate(performance_data, 1):
            table.cell(row, 0).text = overlap
            table.cell(row, 1).text = sc
            table.cell(row, 2).text = mc
            table.cell(row, 3).text = improvement
        
        # 設定表格格式
        for row in range(5):
            for col in range(4):
                cell = table.cell(row, col)
                self.set_mixed_font_style(cell.text_frame, font_size=11)
                if row == 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

    def create_title_slide(self, prs):
        """創建標題投影片"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "LEO 衛星網路多連線換手演算法"
        self.set_mixed_font_style(title.text_frame, font_size=28)
        
        subtitle_content = """Multi-Connectivity Handover Algorithm for LEO Satellite Networks

🎯 演算法核心目標：
• 降低換手中斷率，提升服務連續性
• 利用多重覆蓋區域實現軟切換
• 基於位置預測的智慧觸發機制

📊 主要技術創新：
• Packet Duplication + Selection Combining
• Conditional Handover (CHO) 整合
• Master/Secondary Node 雙連線架構

Based on IEEE Paper: "Enhancing Handover Performance in LEO Satellite Networks"
Mohammed Al-Ansi et al., 2024"""
        
        subtitle.text = subtitle_content
        self.set_mixed_font_style(subtitle.text_frame, font_size=16)

    def create_problem_definition_slide(self, prs):
        """創建問題定義投影片 - 使用分頁控制"""
        
        content = """🛰️ LEO 衛星特性帶來的挑戰：

📡 高速移動特性：
• 軌道速度：7.56 km/s (相比 GEO 靜止不動)
• 覆蓋時間短：~6.6 秒 (50km 波束直徑)
• 頻繁切換需求：每 7 秒需要換手決策

🚫 Single Connectivity Handover (SC-HO) 問題：
• Hard Handover：舊連結先中斷，再建立新連結
• 中斷時間：~50-100ms 不可接受延遲
• 封包遺失：中斷期間資料無法傳輸
• 訊號不穩定：邊界區域連結品質差

⚡ 多重覆蓋區域的機會：
• LEO 星座設計：多顆衛星同時覆蓋
• Overlapping Coverage：30-50% 區域重疊
• 雙連線潛力：同時維持 2+ 衛星連結
• 實現 Soft Handover 的技術基礎"""
        
        # 檢查內容長度並分頁
        content_parts = self.split_content_for_slides(content, self.max_lines_per_slide)
        
        for i, part in enumerate(content_parts):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            if len(content_parts) > 1:
                title_text = f"LEO 衛星網路切換問題分析 ({i+1}/{len(content_parts)})"
            else:
                title_text = "LEO 衛星網路切換問題分析"
            
            slide.shapes.title.text = title_text
            self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
            
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = part
                self.set_mixed_font_style(slide.placeholders[1].text_frame, font_size=15)
            
            # 在第一張投影片添加輔助圖片
            if i == 0:
                self.add_supporting_image(slide, "page_2_img_1.png", "LEO 衛星覆蓋場景")

    def create_algorithm_overview_slide(self, prs):
        """創建演算法概覽投影片 - 使用分頁控制"""
        
        content = """🧠 核心設計理念：
「在多重覆蓋區域同時維持雙重連線，實現無縫軟切換」

🔧 四大技術支柱：

1️⃣ 雙連線架構 (Dual Connectivity)：
   • Master Node (MN)：主要服務衛星，負責控制信令
   • Secondary Node (SN)：輔助衛星，提供數據冗餘
   • 同時連接：避免單點失效風險

2️⃣ 封包複製機制 (Packet Duplication)：
   • 關鍵數據同時從 MN 和 SN 傳輸
   • UE 端接收雙重數據流
   • 選擇式結合：選擇品質最佳數據

3️⃣ 智慧觸發條件 (Conditional Handover)：
   • Location-based 預測：基於軌道位置計算
   • SINR Threshold：信號品質動態門檻
   • 雙重判斷：避免 Ping-pong 效應

4️⃣ 無縫切換機制 (Make-before-break)：
   • 先建立新連結，再釋放舊連結
   • 零中斷時間：連續性服務保證
   • 漸進式切換：平滑過渡避免抖動"""
        
        # 檢查內容長度並分頁
        content_parts = self.split_content_for_slides(content, self.max_lines_per_slide)
        
        for i, part in enumerate(content_parts):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            if len(content_parts) > 1:
                title_text = f"MC-HO 演算法設計理念與核心架構 ({i+1}/{len(content_parts)})"
            else:
                title_text = "MC-HO 演算法設計理念與核心架構"
            
            slide.shapes.title.text = title_text
            self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
            
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = part
                self.set_mixed_font_style(slide.placeholders[1].text_frame, font_size=14)
   • 確保高可靠性資料傳輸

3️⃣ 選擇性合併 (Selection Combining)：
   • 實時監測 MN 和 SN 的 SINR
   • 動態選擇最佳信號品質路徑
   • 最大化接收信號品質

4️⃣ 條件式切換觸發 (Conditional Handover)：
   • 位置基礎觸發：dTSAT(t) ≤ Rb - doffset
   • 預配置機制：提前準備切換參數
   • 降低切換延遲和失效率"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

    def create_algorithm_flowchart_slide(self, prs):
        """創建演算法流程圖投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "MC-HO 演算法執行流程詳解"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        content = """🔄 完整演算法執行流程：

階段一：初始化與監測
├─ UE 連接服務衛星 (SSAT) 作為 Master Node
├─ 持續監測位置：利用 GNSS + 星曆數據
├─ 計算與候選衛星距離：dTSAT(t), dSSAT(t)
└─ 觸發條件檢查：進入多重覆蓋區域？

階段二：條件式觸發 (CHO Activation)
├─ 位置條件：dTSAT(t) ≤ Rb - doffset AND dSSAT(t) ≤ Rb - doffset
├─ 預配置檢查：目標衛星資源是否可用？
├─ 信號品質評估：候選衛星 SINR 是否足夠？
└─ 觸發決策：滿足條件則啟動 MC 架構

階段三：雙連線建立 (Dual Connectivity Setup)
├─ Random Access：與目標衛星建立連線
├─ SN 添加：目標衛星成為 Secondary Node
├─ 封包複製啟動：SSAT 開始數據複製到 TSAT
└─ Selection Combining：UE 開始雙路徑接收

階段四：路徑切換與釋放 (Path Switch & Release)
├─ TSAT 成為新 MN：路徑切換信令
├─ AMF 通知：Bearer Modification 執行
├─ 舊連線釋放：SSAT 連線安全釋放
└─ 準備下次切換：監測下個候選衛星"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=12)

    def create_algorithm_comparison_slide(self, prs):
        """創建演算法比較投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "SC-HO vs MC-HO 演算法技術比較"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        # 添加比較表格
        self.create_algorithm_comparison_table(slide)

    def create_technical_implementation_slide(self, prs):
        """創建技術實現細節投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "MC-HO 關鍵技術實現細節"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        content = """⚙️ 核心演算法實現：

🔧 位置基礎觸發演算法：
```
function checkHandoverTrigger(ue_position, ssat_position, tsat_position):
    d_ssat = calculateDistance(ue_position, ssat_position)
    d_tsat = calculateDistance(ue_position, tsat_position)
    
    if (d_tsat <= Rb - d_offset) AND (d_ssat <= Rb - d_offset):
        return TRIGGER_MC_HO
    else:
        return MAINTAIN_CURRENT
```

📡 封包複製與選擇合併：
```
function selectionCombining(mn_sinr, sn_sinr, mn_data, sn_data):
    if mn_sinr > sn_sinr:
        return mn_data  // 選擇主節點數據
    else:
        return sn_data  // 選擇輔助節點數據
```

🎯 最佳化參數設定：
• 距離偏移量 (d_offset)：1-5 km，影響觸發時機
• 波束半徑 (Rb)：25 km (50km 直徑波束)
• SINR 門檻：-8 dB (連結失效判定)
• 觸發評估週期：0.5 秒 (即時響應)

🔄 狀態機管理：
IDLE → CHO_TRIGGERED → DUAL_CONNECTED → PATH_SWITCHED → RELEASE"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=12)

    def create_performance_analysis_slide(self, prs):
        """創建演算法性能分析投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "MC-HO 演算法性能分析與量化驗證"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
        
        content = """📊 關鍵性能指標分析：

🎯 換手次數改善 (核心指標)：
• 20% 重疊：SC-HO 185 → MC-HO 145 (改善 21.6%)
• 30% 重疊：SC-HO 212 → MC-HO 129 (改善 39.2%)  
• 40% 重疊：SC-HO 247 → MC-HO 130 (改善 47.4%)

🔗 連結可靠性提升：
• RLF 減少：40% 重疊時從 532 降至 410 (改善 23%)
• Selection Combining 效果：平均 SINR 提升 3-5 dB
• 服務中斷時間：從 ~500ms 降至 <50ms

⚡ 演算法複雜度分析：
• 計算複雜度：O(n) - n 為候選衛星數量 (通常 ≤ 3)
• 記憶體開銷：+15% (維持雙連線狀態)
• 信令開銷：+25% (CHO 預配置 + 封包複製)
• 能耗影響：+8% (雙路徑接收)

📈 整體效益評估：
✅ 大幅降低切換頻率 → 提升用戶體驗
✅ 顯著改善可靠性 → 適合關鍵應用  
✅ 合理的資源開銷 → 工程可實現性高"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)
        
        # 添加性能數據表格
        self.create_performance_metrics_table(slide)
        
        # 添加輔助圖片：性能比較圖
        self.add_supporting_image(slide, "page_4_img_1.png", "換手次數比較")

    def create_practical_considerations_slide(self, prs):
        """創建實際應用考量投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "演算法實際部署與工程考量"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        content = """🏗️ 實際部署挑戰與解決方案：

🛰️ 衛星端實現考量：
• Inter-Satellite Link (ISL) 需求：支援 MN-SN 協調
• 計算資源分配：封包複製需要額外處理能力
• 同步要求：雙路徑傳輸的時序同步
• 負載平衡：避免熱點衛星過載

📱 用戶設備 (UE) 要求：
• 雙射頻前端：同時接收多衛星信號
• 協議棧支援：3GPP NTN Dual Connectivity
• 電池續航：多路徑處理的能耗管理
• 天線設計：全向或多波束天線陣列

🌐 網路整合挑戰：
• 核心網改造：支援 MC 架構的 AMF/UPF
• QoS 保證：雙路徑的服務品質協調
• 計費機制：多連線的費用分攤策略
• 互操作性：與現有 4G/5G 網路融合

⚙️ 最佳化策略：
• 動態參數調整：根據網路負載自適應
• 預測性切換：結合軌道預測和用戶行為
• 能效優化：智慧關閉非必要的雙連線
• 故障恢復：連線失效時的快速回退機制"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

    def create_future_work_slide(self, prs):
        """創建未來發展投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "演算法改進方向與未來研究"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        content = """🔮 演算法進階改進方向：

🧠 AI 驅動的智慧切換：
• 機器學習預測：基於歷史數據預測最佳切換時機
• 強化學習優化：Q-learning 動態調整切換參數
• 用戶行為分析：個人化的切換策略
• 網路狀態感知：根據負載動態調整演算法

🌍 Multi-Connectivity 擴展：
• 三重連線：MN + 2×SN 提升可靠性
• 異質網路整合：LEO + GEO + 地面 5G
• 跨軌道平面切換：不同軌道星座間協調
• 動態連線數調整：根據業務需求彈性配置

📡 演算法技術創新：
• 預測性封包複製：只複製關鍵數據包
• 分層 QoS 管理：不同優先級的差異化處理
• 邊緣計算整合：衛星邊緣節點參與決策
• 量子通信準備：未來量子衛星網路支援

⚡ 性能極致優化：
• 毫秒級切換：<10ms 的超低延遲切換
• 零中斷服務：完全無感知的切換體驗
• 能效比優化：降低 50% 的額外能耗
• 大規模部署：支援百萬級用戶同時切換"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

    def create_conclusion_slide(self, prs):
        """創建結論投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "MC-HO 演算法總結與核心貢獻"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        content = """🎯 MC-HO 演算法核心貢獻：

💡 理論創新：
• 首次將 Multi-Connectivity 系統性應用於 LEO 衛星切換
• 位置基礎 CHO 與雙連線架構的有機整合
• 軟切換理論在高動態衛星環境的成功實現

🔧 技術突破：
• Packet Duplication + Selection Combining 機制
• 四階段演算法流程的完整設計與實現
• 複雜度與性能的最佳平衡點

📊 量化成果：
• 換手次數減少：最高達 47.4% (40% 重疊場景)
• 連結失效降低：改善 23% 的可靠性
• 適中的資源開銷：+15% 記憶體，+25% 信令

🌟 實用價值：
• 直接提升 LEO 衛星網路用戶體驗
• 為 6G 衛星通信提供關鍵技術基礎
• 工程可實現性強，具備商用潛力

🚀 影響意義：
MC-HO 演算法為 LEO 衛星網路的大規模商用部署
掃除了關鍵的技術障礙，推動了衛星通信
從輔助角色向主流通信技術的歷史性轉變"""
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

    def generate_algorithm_focused_presentation(self, output_filename="LEO衛星MC-HO演算法深度解析.pptx"):
        """生成以演算法為核心的簡報"""
        
        print("🚀 開始生成以演算法為核心的 PowerPoint 簡報...")
        print("🎯 重點：MC-HO 演算法深度講解，圖表僅作輔助說明")
        print("="*60)
        
        # 1. 載入模板
        prs = self.load_template()
        
        # 2. 按演算法邏輯創建投影片
        print("📄 創建投影片內容...")
        
        self.create_title_slide(prs)
        print("✅ 標題投影片")
        
        self.create_problem_definition_slide(prs)
        print("✅ 問題定義投影片")
        
        self.create_algorithm_overview_slide(prs)
        print("✅ 演算法概覽投影片")
        
        self.create_algorithm_flowchart_slide(prs)
        print("✅ 演算法流程投影片")
        
        self.create_algorithm_comparison_slide(prs)
        print("✅ 演算法比較投影片 (含表格)")
        
        self.create_technical_implementation_slide(prs)
        print("✅ 技術實現投影片")
        
        self.create_performance_analysis_slide(prs)
        print("✅ 性能分析投影片 (含表格+圖)")
        
        self.create_practical_considerations_slide(prs)
        print("✅ 實際應用考量投影片")
        
        self.create_future_work_slide(prs)
        print("✅ 未來發展投影片")
        
        self.create_conclusion_slide(prs)
        print("✅ 結論投影片")
        
        # 3. 儲存簡報
        output_path = f"../../doc/{output_filename}"
        try:
            prs.save(output_path)
            print(f"\n🎉 演算法核心簡報生成成功！")
            print(f"📁 輸出檔案: {output_path}")
            print(f"📊 總投影片數: {len(prs.slides)}")
            
            # 生成報告
            self.generate_creation_report(len(prs.slides))
            
        except Exception as e:
            print(f"❌ 簡報儲存失敗: {e}")

    def generate_creation_report(self, total_slides):
        """生成創建報告"""
        
        report = f"""# 📊 以演算法為核心的簡報創建報告

## 🎯 設計理念
**演算法深度講解為主，圖表輔助說明為輔**

### 核心設計原則：
- 🧠 **演算法優先**: 每張投影片都圍繞 MC-HO 演算法核心內容
- 📊 **圖表輔助**: 圖表僅用於支援演算法說明，非主角
- 📋 **表格整合**: 利用 PowerPoint 表格功能呈現比較數據
- 🔄 **邏輯完整**: 從問題→解決方案→實現→驗證的完整流程

## 📋 簡報結構 ({total_slides} 張投影片)

1. **標題投影片**: 演算法核心目標與技術創新點
2. **問題定義**: LEO 衛星切換挑戰，為演算法設計奠定基礎
3. **演算法概覽**: MC-HO 四大技術支柱詳解
4. **演算法流程**: 四階段執行流程的深度分析
5. **演算法比較**: SC-HO vs MC-HO 技術對比表格
6. **技術實現**: 核心演算法程式碼與參數設定
7. **性能分析**: 量化驗證結果與效益評估 (含表格)
8. **實際應用**: 部署挑戰與工程考量
9. **未來發展**: 演算法改進方向與研究趨勢
10. **結論總結**: 核心貢獻與實用價值

## 🔧 技術特色

### ✅ 演算法深度講解
- 詳細的四階段執行流程
- 具體的程式碼實現片段
- 參數設定與最佳化策略
- 複雜度與性能權衡分析

### ✅ 表格數據呈現
- SC-HO vs MC-HO 技術比較表格
- 量化性能改善數據表格
- 清晰的視覺化數據對比

### ✅ 圖表輔助說明
- 系統覆蓋場景圖 (建立基礎概念)
- 性能比較圖 (驗證演算法效果)
- 圖片位置：右上角小尺寸，不搶奪演算法內容焦點

### ✅ 完整技術邏輯
- 問題定義 → 解決方案設計
- 演算法流程 → 技術實現細節  
- 性能驗證 → 實際部署考量
- 當前成果 → 未來發展方向

## 📈 相比圖表導向版本的改進

| 對比項目 | 圖表導向版本 | 演算法核心版本 |
|---------|-------------|---------------|
| **核心重點** | 圖表展示 | **演算法講解** |
| **內容組織** | 圍繞圖表 | **圍繞演算法邏輯** |
| **技術深度** | 表面描述 | **實現細節+程式碼** |
| **教學價值** | 視覺化展示 | **深度理解算法原理** |
| **實用性** | 了解概念 | **掌握實現方法** |

## ✅ 達成目標確認

- ✅ **演算法為主**: 每張投影片都深入講解 MC-HO 演算法
- ✅ **圖表為輔**: 圖片僅作輔助，不搶奪主要內容焦點  
- ✅ **表格呈現**: 充分利用 PowerPoint 表格功能展示數據
- ✅ **邏輯完整**: 完整的演算法學習和理解路徑
- ✅ **技術實用**: 包含具體實現細節和部署考量

創建時間: {datetime.now().strftime('%Y-%m-%d %H:%M')}
"""

        with open("../../doc/演算法核心簡報創建報告.md", 'w', encoding='utf-8') as f:
            f.write(report)
        
        print("📝 創建報告已儲存至: ../../doc/演算法核心簡報創建報告.md")

def main():
    """主程式"""
    generator = AlgorithmFocusedPowerPointGenerator()
    generator.generate_algorithm_focused_presentation()

if __name__ == "__main__":
    main()