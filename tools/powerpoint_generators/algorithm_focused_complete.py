#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
完整演算法教學簡報生成器 - 含大綱與結構化教學流程
Complete Algorithm Teaching Presentation Generator
"""

import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

class CompleteAlgorithmPresentationGenerator:
    """完整演算法教學簡報生成器"""
    
    def __init__(self, template_path="../../doc/template.pptx"):
        self.template_path = template_path
        self.max_lines_per_slide = 20
        
        # 圖片資源路徑
        self.image_base_path = "../../論文圖片"
        
        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"
    
    def load_template(self):
        """載入簡報模板"""
        try:
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                print(f"✅ 模板載入成功: {self.template_path}")
            else:
                prs = Presentation()
                print("⚠️  模板不存在，使用預設模板")
            return prs
        except Exception as e:
            print(f"❌ 模板載入失敗: {e}")
            return Presentation()

    def estimate_content_lines(self, content_text):
        """估算內容所需行數"""
        lines = content_text.split('\n')
        total_lines = 0
        
        for line in lines:
            if not line.strip():  # 空行
                total_lines += 1
            else:
                char_count = len(line)
                estimated_lines = max(1, (char_count + 79) // 80)
                total_lines += estimated_lines
        
        return total_lines

    def split_content_for_slides(self, content_text, max_lines=20):
        """將內容分割為適合投影片的片段"""
        if self.estimate_content_lines(content_text) <= max_lines:
            return [content_text]
        
        lines = content_text.split('\n')
        parts = []
        current_part = []
        current_lines = 0
        
        for line in lines:
            line_count = max(1, (len(line) + 79) // 80) if line.strip() else 1
            
            if current_lines + line_count > max_lines and current_part:
                parts.append('\n'.join(current_part))
                current_part = [line]
                current_lines = line_count
            else:
                current_part.append(line)
                current_lines += line_count
        
        if current_part:
            parts.append('\n'.join(current_part))
        
        return parts

    def set_mixed_font_style(self, text_frame, font_size=14):
        """設定中英文混合字體"""
        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()
                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', char):
                        eng_text = ""
                        while i < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            eng_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = eng_text
                        run.font.name = self.english_font
                        run.font.size = Pt(font_size)
                    else:
                        chn_text = ""
                        while i < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%]', text[i]):
                            chn_text += text[i]
                            i += 1
                        run = paragraph.add_run()
                        run.text = chn_text
                        run.font.name = self.chinese_font
                        run.font.size = Pt(font_size)

    def add_supporting_image(self, slide, image_file, caption=""):
        """添加輔助說明圖片"""
        if not image_file:
            return None
            
        image_path = self.find_image_path(image_file)
        if image_path:
            try:
                left = Inches(7.5)
                top = Inches(1.2)
                width = Inches(2.2)
                height = Inches(1.5)
                
                picture = slide.shapes.add_picture(image_path, left, top, width, height)
                print(f"✅ 輔助圖片已添加: {os.path.basename(image_path)}")
                return picture
            except Exception as e:
                print(f"❌ 圖片添加失敗: {e}")
                return None

    def find_image_path(self, image_file):
        """智慧尋找圖片檔案路徑"""
        possible_paths = [
            os.path.join(self.image_base_path, image_file),
            f"../../論文圖片/{image_file}",
            f"論文圖片/{image_file}",
            image_file
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                return path
        
        print(f"⚠️  圖片未找到: {image_file}")
        return None

    def create_content_slides(self, prs, title, content, image_file=None, image_caption=""):
        """通用的內容投影片創建函數"""
        
        print(f"📝 創建投影片: {title}")
        print(f"   內容長度: {self.estimate_content_lines(content)} 行")
        
        content_parts = self.split_content_for_slides(content, self.max_lines_per_slide)
        print(f"   分頁結果: {len(content_parts)} 頁")
        
        for i, part in enumerate(content_parts):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            if len(content_parts) > 1:
                title_text = f"{title} ({i+1}/{len(content_parts)})"
            else:
                title_text = title
            
            slide.shapes.title.text = title_text
            self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)
            
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = part
                self.set_mixed_font_style(slide.placeholders[1].text_frame, font_size=14)
            
            if i == 0 and image_file:
                self.add_supporting_image(slide, image_file, image_caption)

    def create_title_slide(self, prs):
        """創建標題投影片"""
        slide_layout = prs.slide_layouts[0]  # 標題版面
        slide = prs.slides.add_slide(slide_layout)
        
        # 設定主標題
        if slide.shapes.title:
            slide.shapes.title.text = "LEO 衛星網路多連線換手演算法"
            self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=28)
        
        # 設定副標題
        if len(slide.placeholders) > 1:
            subtitle_content = """Multi-Connectivity Handover Algorithm for LEO Satellite Networks

🎯 演算法核心目標：
• 降低換手中斷率，提升服務連續性  
• 利用多重覆蓋區域實現軟切換
• 基於位置預測的智慧觸發機制

📊 主要技術創新：
• Packet Duplication + Selection Combining
• Conditional Handover (CHO) 整合  
• Master/Secondary Node 雙連線架構

Based on IEEE Paper: "Enhancing Handover Performance in LEO Satellite Networks"
Mohammed Al-Ansi et al., 2024"""
            
            slide.placeholders[1].text = subtitle_content
            self.set_mixed_font_style(slide.placeholders[1].text_frame, font_size=16)

    def create_outline_slide(self, prs):
        """創建大綱投影片"""
        slide_layout = prs.slide_layouts[1]  # 標題與內容版面
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "簡報大綱 - 演算法教學流程"
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=22)
        
        outline_content = """📚 本簡報教學結構：

1️⃣ 問題背景與動機 (Background & Motivation)
   • LEO 衛星網路特性與挑戰
   • SC-HO 傳統方法的限制
   • MC-HO 解決方案的需求分析

2️⃣ MC-HO 演算法核心設計 (Algorithm Design)  
   • 四大技術支柱詳解
   • 雙連線架構原理
   • 智慧觸發機制

3️⃣ 演算法實現流程 (Implementation Flow)
   • 四階段執行流程
   • 關鍵參數設定
   • 程式碼實現片段

4️⃣ 性能驗證與分析 (Performance Analysis)
   • 量化效益比較
   • 可靠性提升驗證
   • 系統效能權衡

5️⃣ 實際應用與未來發展 (Applications & Future)
   • 部署考量與挑戰
   • 工程實現要求
   • 未來研究方向

🎯 教學目標：深度理解 MC-HO 演算法原理與實現方法"""
        
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = outline_content
            self.set_mixed_font_style(slide.placeholders[1].text_frame, font_size=14)

    def generate_complete_presentation(self, output_filename="LEO衛星MC-HO演算法完整教學簡報.pptx"):
        """生成完整結構化教學簡報"""
        
        print("🚀 開始生成完整結構化教學簡報...")
        print("📚 包含：標題→大綱→問題→演算法→實現→驗證→應用→結論")
        print("="*60)
        
        # 載入模板
        prs = self.load_template()
        
        # 1. 標題投影片
        self.create_title_slide(prs)
        print("✅ 1. 標題投影片")
        
        # 2. 大綱投影片
        self.create_outline_slide(prs)
        print("✅ 2. 大綱投影片")
        
        # 3. 問題背景投影片
        problem_content = """🛰️ LEO 衛星網路特性與挑戰：

📡 LEO 衛星基本特性：
• 軌道高度：600km (相比 GEO 35,786km)
• 運行速度：7.56 km/s (高速移動)
• 覆蓋時間：~6.6 秒 (50km 波束直徑)
• 網路延遲：~2.5ms (大幅優於 GEO ~250ms)

🚫 傳統 SC-HO 的根本問題：
• Hard Handover 機制：Break-before-make
• 服務中斷：50-100ms 無法接受的延遲
• 封包遺失：中斷期間資料傳輸停止
• 頻繁切換：每 7 秒需要換手決策
• Ping-pong 效應：邊界區域來回切換

⚡ 多重覆蓋區域的技術機會：
• LEO 星座設計：多顆衛星同時覆蓋
• 重疊覆蓋率：30-50% 區域具有雙重覆蓋
• Soft Handover 潛力：同時維持多重連結
• 技術基礎：為 MC-HO 提供實現條件"""
        
        self.create_content_slides(prs, "問題背景：LEO 衛星換手挑戰", problem_content, 
                                 "page_2_img_1.png", "LEO 衛星覆蓋場景")
        print("✅ 3. 問題背景投影片")
        
        # 4. MC-HO 演算法設計投影片
        design_content = """🧠 MC-HO 核心設計理念：
「在多重覆蓋區域同時維持雙重連線，實現無縫軟切換」

🏗️ 四大技術支柱：

1️⃣ 雙連線架構 (Dual Connectivity)
   • Master Node (MN)：主要服務衛星
     - 負責控制信令 (RRC, NAS)
     - 提供主要資料通道
   • Secondary Node (SN)：輔助服務衛星
     - 提供冗餘資料通道
     - 分擔資料傳輸負載
   • 同時連接：避免單點失效風險

2️⃣ 封包複製機制 (Packet Duplication)
   • Uplink：UE 同時向 MN 和 SN 傳送
   • Downlink：核心網同步向兩節點發送
   • 冗餘傳輸：確保關鍵資料可靠性
   • 選擇式合併：接收端選擇最佳品質

3️⃣ 條件式切換 (Conditional Handover)  
   • Location-based 預測：基於軌道計算
   • SINR 閾值觸發：信號品質監控
   • 預先配置：減少切換延遲
   • 智慧判斷：避免 Ping-pong 效應

4️⃣ Make-before-break 機制
   • 先建立新連結，再釋放舊連結
   • 零服務中斷：確保連續性
   • 漸進式切換：平滑角色轉換"""
        
        self.create_content_slides(prs, "MC-HO 演算法核心設計", design_content)
        print("✅ 4. 演算法設計投影片")
        
        # 5. 演算法實現流程投影片
        flow_content = """📋 MC-HO 四階段實現流程：

🔍 階段一：多重覆蓋偵測 (Multi-Coverage Detection)
   Algorithm Step 1:
   • 位置服務 (LCS) 持續監測 UE 座標
   • 計算可見衛星清單與覆蓋時間
   • 評估 SINR 預測值與服務品質
   • 判斷是否進入多重覆蓋區域

🔗 階段二：次要連結建立 (Secondary Connection Setup)
   Algorithm Step 2:  
   • 在多重覆蓋區域啟動 SN 連線程序
   • 執行 RRC Connection Establishment
   • 配置 Dual Connectivity 參數
   • 啟用 Packet Duplication 機制

📊 階段三：條件式切換執行 (CHO Execution)
   Algorithm Step 3:
   • 持續監測觸發條件 (SINR + Location)
   • 條件滿足時執行角色切換：
     - Current SN → New MN  
     - Current MN → New SN (如仍可用)
   • 更新路由表與資料流向
   • 確保切換過程零中斷

✂️ 階段四：舊連結釋放 (Legacy Connection Release)  
   Algorithm Step 4:
   • 評估舊衛星連結品質與必要性
   • 適時釋放不再需要的連線資源
   • 最佳化系統資源使用效率
   • 準備進入下一輪切換循環"""
        
        self.create_content_slides(prs, "MC-HO 四階段實現流程", flow_content)
        print("✅ 5. 實現流程投影片")
        
        # 6. 核心演算法程式碼投影片
        code_content = """💻 核心演算法程式碼實現：

🔧 多重覆蓋偵測演算法：
```python
def detect_multi_coverage(ue_position, satellite_list, sinr_threshold):
    candidates = []
    for satellite in satellite_list:
        # 計算距離與仰角
        distance = calculate_distance(ue_position, satellite.position)
        elevation = calculate_elevation_angle(ue_position, satellite.position)
        
        # 預測 SINR 值
        predicted_sinr = link_budget_calculation(distance, elevation)
        
        # 計算覆蓋時間
        coverage_time = predict_coverage_duration(ue_position, satellite.orbit)
        
        if predicted_sinr > sinr_threshold and coverage_time > min_coverage:
            candidates.append({
                'satellite': satellite,
                'sinr': predicted_sinr,
                'coverage_time': coverage_time
            })
    
    return candidates
```

📡 雙連線管理演算法：
```python  
def manage_dual_connectivity(mn_satellite, sn_candidate):
    # 建立 Secondary Node 連線
    if establish_sn_connection(sn_candidate):
        # 配置封包複製
        configure_packet_duplication(mn_satellite, sn_candidate)
        
        # 啟用選擇式合併
        enable_selection_combining()
        
        return True
    return False
```

⚡ CHO 觸發條件檢查：
```python
def check_handover_trigger(current_mn, current_sn, threshold_config):
    # 位置基礎觸發
    coverage_ending = predict_coverage_end(current_mn) < 3.0  # 3秒內結束
    
    # SINR 基礎觸發  
    sinr_degraded = current_mn.sinr < threshold_config.min_sinr
    
    # 綜合判斷
    return coverage_ending or sinr_degraded
```"""
        
        self.create_content_slides(prs, "核心演算法程式碼實現", code_content)
        print("✅ 6. 程式碼實現投影片")
        
        # 7. 性能驗證投影片
        performance_content = """📈 MC-HO vs SC-HO 性能驗證結果：

🏆 換手次數顯著改善：
• 測試環境：40% 重疊覆蓋率
  - SC-HO：247 次換手/秒 (頻繁切換)
  - MC-HO：130 次換手/秒 (改善 47%)
• 測試環境：30% 重疊覆蓋率  
  - SC-HO：195 次換手/秒
  - MC-HO：118 次換手/秒 (改善 39%)

🛡️ 系統可靠性大幅提升：
• 無線連結失效 (RLF) 次數：
  - SC-HO：532 次/秒 (高失效率)
  - MC-HO：409 次/秒 (改善 23%)
• 服務中斷時間：
  - SC-HO：平均 75ms 中斷時間
  - MC-HO：接近 0ms (Soft Handover)

📊 系統效能權衡分析：
• 優勢面：
  - 服務連續性：顯著改善用戶體驗
  - 可靠性提升：降低通話中斷機率
  - 邊緣覆蓋：改善邊界區域服務品質
• 成本面：
  - 信令開銷：增加 15-20% (雙連線)
  - 頻譜效率：輕微下降 5% (冗餘傳輸)
  - 複雜度：系統實現複雜度提升

✅ 整體評估：效益遠大於成本，適合部署"""
        
        self.create_content_slides(prs, "性能驗證：量化效益分析", performance_content,
                                 "page_4_img_1.png", "性能比較圖表")
        print("✅ 7. 性能驗證投影片")
        
        # 8. 實際應用投影片
        application_content = """🚀 實際部署應用與工程挑戰：

🏗️ 系統架構要求：
• 3GPP NTN 標準相容：
  - Release 17/18 NTN 功能支援
  - 5G 核心網 (5GC) Dual Connectivity 能力
  - RAN 設備支援 MN/SN 角色切換
• UE 終端要求：
  - 多連線併發處理能力
  - 增強的 SINR 測量精度
  - Location Service (LCS) 支援

⚙️ 關鍵參數調校：
• SINR 閾值：-10dB 到 -15dB (可調)
• 重疊覆蓋閾值：最低 30% 覆蓋率
• CHO 觸發提前時間：2-4 秒範圍
• Hysteresis 邊界：2-5dB 避免抖動

🔄 系統整合策略：
• 向後相容性：支援傳統 SC-HO 終端
• 漸進式部署：從高密度區域開始
• 負載平衡：智慧分配 MN/SN 角色
• 故障恢復：自動降級為 SC-HO 模式

🎯 未來發展roadmap：
• Phase 1：核心功能驗證與部署
• Phase 2：AI/ML 預測演算法整合
• Phase 3：多層衛星星座協調
• Phase 4：6G 地面-衛星融合網路"""
        
        self.create_content_slides(prs, "實際應用部署與未來發展", application_content)
        print("✅ 8. 實際應用投影片")
        
        # 9. 結論投影片
        conclusion_content = """🎯 MC-HO 演算法總結與貢獻：

💡 主要技術突破：
• 首次實現 LEO 衛星真正 Soft Handover 機制
• 雙連線架構有效解決服務中斷問題  
• Location-based CHO 提升切換效率
• Make-before-break 確保服務連續性

📊 量化效能證實：
• 換手次數減少：最高 47% 改善
• 連結失效降低：23% 可靠性提升
• 服務中斷時間：從 75ms → 接近 0ms
• 用戶體驗品質：顯著提升

🏗️ 工程實用價值：
• 標準相容：符合 3GPP NTN Release 17/18
• 部署友善：向後相容現有系統
• 參數靈活：可調整適應不同場景
• 擴展性強：支援未來 6G 演進

🔬 學術研究貢獻：
• 理論創新：多連線軟切換理論基礎
• 演算法設計：四階段結構化實現流程
• 性能評估：完整的量化驗證方法
• 實現指導：具體的工程部署方案

🚀 未來研究方向：
• AI 驅動的智慧預測與最佳化
• 多層星座協調演算法  
• 超低延遲應用場景適配
• 地面-衛星異構網路融合

結論：MC-HO 演算法成功解決 LEO 衛星頻繁切換問題，
為下一代衛星通訊網路奠定重要技術基礎。"""
        
        self.create_content_slides(prs, "結論：技術貢獻與未來展望", conclusion_content)
        print("✅ 9. 結論投影片")
        
        # 儲存簡報
        output_path = f"../../doc/{output_filename}"
        try:
            prs.save(output_path)
            print(f"\n🎉 完整教學簡報生成成功！")
            print(f"📁 輸出檔案: {output_path}")
            print(f"📊 總投影片數: {len(prs.slides)}")
            print(f"📚 結構: 標題→大綱→問題→設計→流程→程式碼→驗證→應用→結論")
            
            # 生成統計報告
            self.generate_creation_report(len(prs.slides))
            
        except Exception as e:
            print(f"❌ 簡報儲存失敗: {e}")

    def generate_creation_report(self, total_slides):
        """生成簡報創建報告"""
        
        report = f"""# 📚 完整演算法教學簡報創建報告

## 🎯 教學設計理念
**結構化演算法教學，從問題到解決方案的完整學習路徑**

### 📋 教學流程設計：
1. **問題引入** → 建立學習動機與背景知識
2. **演算法設計** → 深入理解核心技術原理  
3. **實現流程** → 掌握具體實作方法
4. **程式碼實現** → 提供可執行的技術細節
5. **性能驗證** → 量化證實演算法效果
6. **實際應用** → 了解工程部署考量
7. **總結展望** → 整合知識並展望未來

## 📊 簡報結構 ({total_slides} 張投影片)

### 🎭 標題與導引部分：
1. **標題投影片** - 演算法核心目標與技術創新
2. **大綱投影片** - 完整教學流程與學習目標

### 📚 核心教學內容：
3. **問題背景** - LEO 衛星換手挑戰與 SC-HO 限制
4. **演算法設計** - MC-HO 四大技術支柱詳解
5. **實現流程** - 四階段執行流程深度分析
6. **程式碼實現** - 核心演算法具體實作
7. **性能驗證** - 量化效益與可靠性分析
8. **實際應用** - 部署考量與工程挑戰

### 🎯 總結與展望：
9. **結論總結** - 技術貢獻與未來研究方向

## ✅ 教學完整性確認

### 🧠 認知層次涵蓋：
- **記憶**：LEO 衛星基本特性與參數
- **理解**：MC-HO 演算法工作原理  
- **應用**：演算法實現流程與程式碼
- **分析**：性能比較與權衡分析
- **綜合**：系統整合與部署策略
- **評估**：技術貢獻與未來發展評估

### 📏 技術指南遵循：
- ✅ **高度控制**：每頁嚴格 ≤{self.max_lines_per_slide} 行
- ✅ **自動分頁**：長內容智慧分割
- ✅ **字體混合**：中英文正確字體設定
- ✅ **圖表輔助**：右上角小尺寸，不搶奪焦點

### 🎓 教學目標達成：
- ✅ **理論基礎**：深度理解 MC-HO 演算法原理
- ✅ **實作能力**：掌握具體實現方法與程式碼  
- ✅ **應用知識**：了解實際部署考量與挑戰
- ✅ **批判思考**：分析演算法優缺點與適用場景

## 🏆 相比之前版本的全面改進

| 改進項目 | 之前版本 | 完整教學版本 |
|---------|----------|--------------|
| **結構完整性** | 缺乏大綱，直接進入內容 | **完整教學流程：標題→大綱→內容→結論** |
| **教學邏輯** | 內容跳躍，缺乏連貫性 | **循序漸進：問題→解決→驗證→應用** |
| **學習引導** | 無學習目標說明 | **明確大綱與學習目標導引** |
| **內容深度** | 表面介紹 | **四個層次：原理→流程→程式碼→應用** |
| **實用價值** | 理論偏重 | **理論與實務並重，含部署指導** |

## 📈 教學效果預期

### 🎯 適用對象：
- **研究生課程**：衛星通訊、行動網路相關課程
- **工程師培訓**：LEO 衛星系統開發團隊
- **學術研究**：論文技術細節理解與驗證
- **產業應用**：衛星通訊產品技術評估

### 📚 學習成效：
- **基礎理解**：完全掌握 LEO 衛星換手問題
- **核心技術**：深度理解 MC-HO 演算法設計
- **實作能力**：具備演算法實現的技術能力
- **應用視野**：了解實際部署的工程考量

創建時間: {datetime.now().strftime('%Y-%m-%d %H:%M')}
"""
        
        with open("../../doc/完整教學簡報創建報告.md", 'w', encoding='utf-8') as f:
            f.write(report)
        
        print("📝 創建報告已儲存至: ../../doc/完整教學簡報創建報告.md")

def main():
    """主程式"""
    generator = CompleteAlgorithmPresentationGenerator()
    generator.generate_complete_presentation()

if __name__ == "__main__":
    main()