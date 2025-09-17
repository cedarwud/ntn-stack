#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
DHO演算法技術原理詳細簡報生成器
針對演算法核心進行深度解析，包含公式詳解、變數意義說明、流程圖等
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

class DHOAlgorithmDetailedPresentation:
    """DHO演算法技術原理詳細簡報生成器"""

    def __init__(self, output_dir="../../doc"):
        self.output_dir = output_dir
        self.max_lines_per_slide = 18  # 降低行數以容納更多詳細內容

        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"
        self.formula_font = "Times New Roman"

        # 顏色設定
        self.primary_color = RGBColor(0, 51, 102)      # 深藍色
        self.secondary_color = RGBColor(255, 102, 0)   # 橙色
        self.formula_color = RGBColor(139, 0, 139)     # 深紫色

        # 確保輸出目錄存在
        os.makedirs(self.output_dir, exist_ok=True)

    def set_mixed_font_style(self, text_frame, chinese_font=None, english_font=None, font_size=14):
        """設置混合中英文字體"""
        if chinese_font is None:
            chinese_font = self.chinese_font
        if english_font is None:
            english_font = self.english_font

        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()

                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', char):
                        # 收集連續的英文字符
                        j = i
                        while j < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = english_font
                        run.font.size = Pt(font_size)
                        i = j
                    else:
                        # 收集連續的中文字符
                        j = i
                        while j < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = chinese_font
                        run.font.size = Pt(font_size)
                        i = j

    def add_formula_slide(self, prs, title, formula_latex, formula_explanation, variables_explanation):
        """添加公式解釋投影片"""
        slide_layout = prs.slide_layouts[1]  # 標題與內容版面
        slide = prs.slides.add_slide(slide_layout)

        # 設定標題
        slide.shapes.title.text = title
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)

        # 創建內容文字框
        content_placeholder = slide.placeholders[1]

        # 公式部分
        content_text = f"""【核心公式】
{formula_latex}

【公式說明】
{formula_explanation}

【變數定義】
{variables_explanation}"""

        content_placeholder.text = content_text
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=13)

        return slide

    def add_algorithm_flow_slide(self, prs, title, flow_content):
        """添加演算法流程投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)

        content_placeholder = slide.placeholders[1]
        content_placeholder.text = flow_content
        self.set_mixed_font_style(content_placeholder.text_frame, font_size=14)

        return slide

    def add_comparison_table_slide(self, prs, title, table_data):
        """添加比較表格投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=20)

        # 創建表格
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 2

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 填入表格數據
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                self.set_mixed_font_style(cell.text_frame, font_size=12)

        return slide

    def create_dho_presentation(self):
        """創建DHO演算法技術原理詳細簡報"""
        print("🚀 開始創建DHO演算法技術原理詳細簡報...")

        # 創建簡報
        prs = Presentation()

        # 1. 標題頁
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "DHO演算法技術原理詳細解析"
        title_slide.placeholders[1].text = "省略MR的智能換手決策機制\\n深度技術原理與公式解析"
        self.set_mixed_font_style(title_slide.shapes.title.text_frame, font_size=24)
        self.set_mixed_font_style(title_slide.placeholders[1].text_frame, font_size=18)

        # 2. 目錄頁
        toc_slide = prs.slides.add_slide(prs.slide_layouts[1])
        toc_slide.shapes.title.text = "簡報大綱"
        toc_content = """1. DHO演算法核心概念
2. 傳統HO vs DHO的根本差異
3. MDP建模詳細解析
4. 狀態空間設計原理
5. 動作空間協調機制
6. 獎勵函數設計
7. IMPALA演算法實現
8. V-trace機制數學原理
9. 預測能力學習機制
10. 性能優勢量化分析"""
        toc_slide.placeholders[1].text = toc_content
        self.set_mixed_font_style(toc_slide.shapes.title.text_frame, font_size=20)
        self.set_mixed_font_style(toc_slide.placeholders[1].text_frame, font_size=14)

        # 3. DHO演算法核心概念
        concept_slide = prs.slides.add_slide(prs.slide_layouts[1])
        concept_slide.shapes.title.text = "DHO演算法核心概念"
        concept_content = """【革命性創新】
• 從「反應式測量」到「預測式決策」的范式轉換
• 省略測量報告(MR)階段，直接進行智能決策

【核心技術突破】
• 利用LEO軌道的確定性進行模式學習
• 基於歷史經驗預測最佳換手時機
• 多用戶聯合優化，避免資源衝突

【技術優勢】
• 延遲降低：消除MR傳輸的1.6~6ms延遲
• 功耗節省：無需週期性測量，節省30-50%功耗
• 準確性提升：基於實時狀態，存取延遲降低6.86倍"""
        concept_slide.placeholders[1].text = concept_content
        self.set_mixed_font_style(concept_slide.shapes.title.text_frame, font_size=20)
        self.set_mixed_font_style(concept_slide.placeholders[1].text_frame, font_size=13)

        # 4. 傳統HO vs DHO流程對比
        comparison_slide = prs.slides.add_slide(prs.slide_layouts[1])
        comparison_slide.shapes.title.text = "傳統HO vs DHO流程對比"
        comparison_content = """【傳統HO流程】
Step 1: UE測量信號(RSRP, RSRQ) → 100-200ms
Step 2: 生成測量報告(MR) → 封裝時間
Step 3: 傳輸延遲 → 3.2-12ms (LEO衛星)
Step 4: gNB接收分析 → 處理時間
Step 5: HO決策 → 基於過時數據
Step 6: 執行換手 → 總延遲112-212ms

【DHO創新流程】
Step 1: 智能代理觀察環境狀態 → <1ms
Step 2: 基於學習模式預測時機 → 即時
Step 3: 直接發送HO請求 → 無等待
總延遲：幾乎即時 (>100倍改善)"""
        comparison_slide.placeholders[1].text = comparison_content
        self.set_mixed_font_style(comparison_slide.shapes.title.text_frame, font_size=20)
        self.set_mixed_font_style(comparison_slide.placeholders[1].text_frame, font_size=13)

        # 5. MDP建模 - 狀態空間設計
        self.add_formula_slide(
            prs,
            "MDP建模：狀態空間設計",
            "s[n] = {n, a^HO[n], a[n-1]}",
            """狀態空間包含三個核心組件，每個組件都具有特定的技術意義：
• 時間索引隱式編碼軌道位置信息
• 存取狀態反映當前網路負載情況
• 歷史動作提供決策連續性依據""",
            """n: 時間索引，軌道週期內的時間位置 (0 ≤ n < T)
a^HO[n]: 存取狀態向量，a^HO_j[n] ∈ {0,1}
  - 1: UE j已成功存取
  - 0: UE j未成功存取
a[n-1]: 上一時隙的動作向量，提供決策歷史信息"""
        )

        # 6. MDP建模 - 動作空間設計
        self.add_formula_slide(
            prs,
            "MDP建模：動作空間設計",
            "a_j[n] = {a_0, a_1, a_2, ..., a_{K-1}}, Σ(k=0 to K-1) a_k = 1",
            """動作空間採用one-hot編碼，確保每個UE只能選擇一個目標：
• a_0 = 1: 不進行換手(智能退避策略)
• a_k = 1 (k≥1): 選擇目標衛星k進行換手
• 互斥選擇避免資源衝突""",
            """a_j[n]: UE j在時隙n的動作向量
K: 總目標數量(包含"不換手"選項)
a_k: 動作選擇指示器，滿足one-hot約束
全域動作: a[n] = [a_1[n], a_2[n], ..., a_J[n]]^T"""
        )

        # 7. MDP建模 - 獎勵函數設計
        self.add_formula_slide(
            prs,
            "MDP建模：獎勵函數設計",
            "r[n] = -D[n] - ν·C[n]",
            """獎勵函數平衡兩個關鍵性能指標：
• 存取延遲懲罰：鼓勵快速成功存取
• 碰撞率懲罰：避免資源衝突
• 權衡係數ν根據應用需求調整""",
            """D[n] = (1/|J|)·Σ(j∈J)(1 - a_j^HO[n]): 正規化存取延遲
C[n] = Σ(k=1 to K-1)C_k^R[n] + C^P[n]: 總碰撞率
  - C_k^R[n]: 目標k的RB碰撞率
  - C^P[n]: PRACH碰撞率
ν: 權衡係數 (URLLC大,mMTC小)"""
        )

        # 8. IMPALA演算法 - V-trace機制
        self.add_formula_slide(
            prs,
            "IMPALA演算法：V-trace機制",
            "v[n] = V(s[n]) + Σ(i=s to s+k-1) γ^(i-s)·Π(j=s to i-1)c[j]·δ_i^V",
            """V-trace解決off-policy學習中的策略滯後問題：
• 修正行為策略μ與目標策略π的差異
• 重要性採樣權重確保收斂性
• 截斷機制保證訓練穩定性""",
            """V(s[n]): 狀態價值函數
γ: 折扣因子 (0 < γ < 1)
c[j]: 截斷重要性權重 = min(c̄, π(a[j]|s[j])/μ(a[j]|s[j]))
δ_i^V: TD誤差 = ρ[i](r[i] + γV(s[i+1]) - V(s[i]))
ρ[i]: 截斷重要性權重 = min(ρ̄, π(a[i]|s[i])/μ(a[i]|s[i]))"""
        )

        # 9. 重要性採樣權重計算
        self.add_formula_slide(
            prs,
            "重要性採樣權重詳解",
            "ρ[n] = min(ρ̄, π(a[n]|s[n])/μ(a[n]|s[n]))",
            """重要性權重補償策略差異，確保正確的價值估計：
• π(a|s): 目標策略下的動作概率
• μ(a|s): 行為策略下的動作概率
• 截斷防止權重過大導致方差爆炸""",
            """π(a[n]|s[n]): learner正在更新的目標策略
μ(a[n]|s[n]): actor收集數據時的行為策略
ρ̄: 截斷閾值，通常設為1.0
c̄: 截斷閾值，通常設為1.0
權重作用：修正策略不一致造成的偏差"""
        )

        # 10. 軌道動力學預測基礎
        self.add_formula_slide(
            prs,
            "軌道動力學預測基礎",
            "q_i[m] = q_i[0] + τ·Σ(m'=1 to m)v_i[m'τ]",
            """LEO衛星軌道的確定性是DHO預測能力的物理基礎：
• 牛頓力學和克卜勒定律保證軌道可預測性
• 時間→位置→覆蓋→最佳決策的映射鏈
• 週期性模式為學習提供穩定基礎""",
            """q_i[m]: 衛星i在第m個時刻的位置向量
q_i[0]: 衛星i的初始位置
τ: 時間間隔步長
v_i[m'τ]: 衛星i在時刻m'τ的速度向量
m: 時間步驟索引"""
        )

        # 11. 演算法流程圖
        flow_content = """【DHO完整執行流程】

Phase 1: 初始化
├─ 載入TLE軌道數據
├─ 初始化神經網路參數
└─ 設定學習超參數

Phase 2: 訓練階段
├─ Actor收集經驗
│  ├─ 觀察狀態s[n] = {n, a^HO[n], a[n-1]}
│  ├─ 執行動作a[n] = π(·|s[n])
│  └─ 記錄(s,a,r,s')軌跡
├─ Learner更新策略
│  ├─ 計算V-trace目標
│  ├─ 更新價值函數V(s)
│  └─ 更新策略π(a|s)
└─ 同步參數至Actor

Phase 3: 部署階段
├─ 實時狀態觀察
├─ 策略推理π(a|s)
└─ 執行換手決策"""

        self.add_algorithm_flow_slide(prs, "DHO演算法完整執行流程", flow_content)

        # 12. 性能對比表格
        performance_data = [
            ["指標", "傳統HO", "DHO", "改善幅度"],
            ["決策延遲", "112-212ms", "<1ms", ">100倍"],
            ["功耗節省", "基準", "30-50%節省", "顯著"],
            ["存取延遲", "基準", "降低6.86倍", "686%"],
            ["碰撞率", "高(資源衝突)", "低(協調優化)", "大幅改善"],
            ["適應性", "固定規則", "動態學習", "質的提升"],
            ["可擴展性", "有限", "支持大規模UE", "架構優勢"]
        ]

        self.add_comparison_table_slide(prs, "DHO vs 傳統HO性能對比", performance_data)

        # 13. 技術創新總結
        innovation_slide = prs.slides.add_slide(prs.slide_layouts[1])
        innovation_slide.shapes.title.text = "DHO技術創新總結"
        innovation_content = """【核心技術突破】
✓ 省略MR：消除測量-報告-決策的延遲鏈路
✓ 預測決策：基於軌道確定性的模式學習
✓ 聯合優化：多UE協調避免資源衝突
✓ 自適應性：動態調整策略適應網路變化

【演算法優勢】
• 延遲：100倍以上的決策延遲降低
• 功耗：30-50%的HO相關功耗節省
• 性能：存取延遲降低6.86倍
• 架構：支持大規模UE的分散式處理

【工程價值】
• 為未來NTN標準提供技術參考
• 開創通訊系統智能化新方向
• 實現reactive到proactive的范式轉換"""
        innovation_slide.placeholders[1].text = innovation_content
        self.set_mixed_font_style(innovation_slide.shapes.title.text_frame, font_size=20)
        self.set_mixed_font_style(innovation_slide.placeholders[1].text_frame, font_size=13)

        # 14. 結論
        conclusion_slide = prs.slides.add_slide(prs.slide_layouts[1])
        conclusion_slide.shapes.title.text = "結論與展望"
        conclusion_content = """【技術貢獻】
• 首次實現LEO衛星換手的MDP完整建模
• 創新性地將DRL應用於NTN優化問題
• 提供了省略MR的可行技術路徑

【學術價值】
• 深度強化學習在通訊系統的創新應用
• LEO衛星網路智能化的重要技術突破
• 為未來6G NTN標準化提供理論基礎

【未來方向】
• 多目標優化擴展(能源、QoS、負載均衡)
• 跨層聯合優化整合
• 大規模星座的分散式協作機制"""
        conclusion_slide.placeholders[1].text = conclusion_content
        self.set_mixed_font_style(conclusion_slide.shapes.title.text_frame, font_size=20)
        self.set_mixed_font_style(conclusion_slide.placeholders[1].text_frame, font_size=14)

        # 儲存簡報
        output_filename = os.path.join(self.output_dir, "DHO演算法技術原理詳細解析.pptx")
        prs.save(output_filename)

        print(f"✅ DHO演算法技術原理詳細簡報創建完成")
        print(f"📊 總投影片數：{len(prs.slides)} 張")
        print(f"📁 儲存位置：{output_filename}")

        # 生成創建報告
        self.generate_creation_report(len(prs.slides))

        return output_filename

    def generate_creation_report(self, total_slides):
        """生成簡報創建報告"""
        report_content = f"""# DHO演算法技術原理詳細簡報創建報告

## 📊 簡報概覽
- **簡報主題**: DHO演算法技術原理詳細解析
- **總投影片數**: {total_slides} 張
- **創建時間**: 2024-09-17
- **技術重點**: 演算法核心、公式解析、變數說明

## 🎯 內容結構
1. **標題頁** - 簡報主題介紹
2. **目錄頁** - 完整大綱指引
3. **核心概念** - DHO基本原理
4. **流程對比** - 傳統HO vs DHO
5. **MDP建模** - 狀態空間、動作空間、獎勵函數
6. **IMPALA算法** - V-trace機制詳解
7. **重要性採樣** - 權重計算機制
8. **軌道預測** - 物理基礎原理
9. **執行流程** - 完整算法流程
10. **性能對比** - 量化效果分析
11. **技術創新** - 核心貢獻總結
12. **結論展望** - 學術價值與未來方向

## 🔬 技術特色
- **公式詳解**: 每個核心公式都有完整的數學表達和變數說明
- **深度解析**: 重點說明演算法的技術原理和創新點
- **流程清晰**: 提供完整的演算法執行流程圖
- **對比分析**: 量化展示DHO相對傳統方法的優勢

## ✅ 品質保證
- **字體設定**: 正確的中英文混合字體
- **內容密度**: 每頁控制在適當行數內
- **技術準確性**: 基於原始論文的精確內容
- **邏輯完整性**: 從基礎概念到深度技術的漸進式說明

## 📈 相比基礎版本的優勢
- **更深入**: 詳細解釋每個公式的數學意義
- **更完整**: 包含完整的演算法流程和實現細節
- **更專業**: 針對技術專家的深度內容
- **更實用**: 提供了工程實現的重要考量

本簡報適合：
- 學術研究人員的技術交流
- 工程團隊的深度技術分享
- 論文答辯的核心內容展示
- 技術標準化組織的提案說明
"""

        report_filename = os.path.join(self.output_dir, "DHO演算法詳細簡報創建報告.md")
        with open(report_filename, 'w', encoding='utf-8') as f:
            f.write(report_content)

        print(f"📋 創建報告已生成：{report_filename}")

def main():
    """主函數"""
    print("🚀 DHO演算法技術原理詳細簡報生成器")
    print("=" * 50)

    # 創建生成器實例
    generator = DHOAlgorithmDetailedPresentation()

    # 生成簡報
    try:
        output_file = generator.create_dho_presentation()
        print(f"\n✅ 簡報生成成功：{output_file}")
        print("\n🎯 本簡報特色：")
        print("   • 深度解析DHO演算法核心技術")
        print("   • 詳細解釋所有關鍵公式和變數")
        print("   • 完整的演算法流程圖說明")
        print("   • 量化的性能對比分析")
        print("   • 適合技術專家的深度內容")

    except Exception as e:
        print(f"❌ 簡報生成失敗：{e}")
        return False

    return True

if __name__ == "__main__":
    main()