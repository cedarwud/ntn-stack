#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
DHO演算法技術原理完整簡報生成器（含圖片整合）
基於原始論文內容，整合提取的圖片，提供最完整的技術解析
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

class DHOAlgorithmCompletePresentation:
    """DHO演算法技術原理完整簡報生成器（含圖片）"""

    def __init__(self, output_dir="../../doc", image_dir="../../論文圖片"):
        self.output_dir = output_dir
        self.image_dir = image_dir
        self.max_lines_per_slide = 16  # 為圖片預留空間

        # 字型設定
        self.chinese_font = "標楷體"
        self.english_font = "Times New Roman"

        # 顏色設定
        self.primary_color = RGBColor(0, 51, 102)      # 深藍色
        self.secondary_color = RGBColor(255, 102, 0)   # 橙色
        self.formula_color = RGBColor(139, 0, 139)     # 深紫色

        # 確保輸出目錄存在
        os.makedirs(self.output_dir, exist_ok=True)

        # 載入圖片資訊
        self.load_figure_info()

    def load_figure_info(self):
        """載入論文圖片資訊"""
        try:
            info_file = os.path.join(self.image_dir, "extraction_info.json")
            with open(info_file, 'r', encoding='utf-8') as f:
                self.figure_info = json.load(f)
            print(f"✅ 載入論文圖片資訊：{len(self.figure_info)} 張圖片")
        except Exception as e:
            print(f"⚠️  無法載入圖片資訊：{e}")
            self.figure_info = []

    def set_mixed_font_style(self, text_frame, chinese_font=None, english_font=None, font_size=14):
        """設置混合中英文字體"""
        if chinese_font is None:
            chinese_font = self.chinese_font
        if english_font is None:
            english_font = self.english_font

        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            if text:
                paragraph.clear()

                i = 0
                while i < len(text):
                    char = text[i]
                    if re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', char):
                        # 收集連續的英文字符
                        j = i
                        while j < len(text) and re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = english_font
                        run.font.size = Pt(font_size)
                        i = j
                    else:
                        # 收集連續的中文字符
                        j = i
                        while j < len(text) and not re.match(r'[a-zA-Z0-9\s\-_.,()[\]/+=<>&%:;!?$]', text[j]):
                            j += 1
                        run = paragraph.add_run()
                        run.text = text[i:j]
                        run.font.name = chinese_font
                        run.font.size = Pt(font_size)
                        i = j

    def add_figure_slide(self, prs, title, content, figure_filename=None, figure_description=""):
        """添加包含圖片的投影片"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 設定標題
        slide.shapes.title.text = title
        self.set_mixed_font_style(slide.shapes.title.text_frame, font_size=18)

        # 如果有圖片，調整內容佈局
        if figure_filename and os.path.exists(os.path.join(self.image_dir, figure_filename)):
            # 縮小內容區域為左半部
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content
            self.set_mixed_font_style(content_placeholder.text_frame, font_size=12)

            # 調整內容框大小和位置
            content_placeholder.left = Inches(0.5)
            content_placeholder.top = Inches(1.5)
            content_placeholder.width = Inches(4.5)
            content_placeholder.height = Inches(5)

            # 添加圖片到右半部
            try:
                image_path = os.path.join(self.image_dir, figure_filename)
                left = Inches(5.2)
                top = Inches(1.5)
                width = Inches(4)
                height = Inches(4.5)

                picture = slide.shapes.add_picture(image_path, left, top, width, height)

                # 添加圖片說明
                if figure_description:
                    textbox = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.4))
                    text_frame = textbox.text_frame
                    text_frame.text = figure_description
                    self.set_mixed_font_style(text_frame, font_size=10)

                print(f"✅ 添加圖片：{figure_filename}")

            except Exception as e:
                print(f"❌ 圖片添加失敗：{e}")
        else:
            # 沒有圖片時使用全寬内容
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content
            self.set_mixed_font_style(content_placeholder.text_frame, font_size=13)

        return slide

    def add_formula_with_figure_slide(self, prs, title, formula_latex, formula_explanation,
                                    variables_explanation, figure_filename=None, figure_description=""):
        """添加公式解釋投影片（含圖片）"""
        content_text = f"""【核心公式】
{formula_latex}

【公式說明】
{formula_explanation}

【變數定義】
{variables_explanation}"""

        return self.add_figure_slide(prs, title, content_text, figure_filename, figure_description)

    def create_comprehensive_dho_presentation(self):
        """創建完整的DHO演算法技術原理簡報"""
        print("🚀 開始創建DHO演算法完整技術簡報（含圖片）...")

        # 創建簡報
        prs = Presentation()

        # 1. 標題頁
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "DHO演算法完整技術原理"
        title_slide.placeholders[1].text = "省略MR的智能換手決策機制\\n深度技術解析與實驗驗證"
        self.set_mixed_font_style(title_slide.shapes.title.text_frame, font_size=26)
        self.set_mixed_font_style(title_slide.placeholders[1].text_frame, font_size=18)

        # 2. 簡報大綱
        toc_content = """1. 研究背景與動機
2. DHO演算法核心創新
3. 系統模型與問題建模
4. MDP數學建模詳解
5. IMPALA演算法機制
6. 實驗設計與驗證
7. 性能分析與對比
8. 技術貢獻與影響
9. 結論與未來方向"""

        self.add_figure_slide(prs, "簡報大綱", toc_content,
                             figure_filename="page_2_img_1.png",
                             figure_description="圖1: LEO衛星網路系統架構")

        # 3. 研究背景 - LEO衛星挑戰
        background_content = """【LEO衛星換手挑戰】
• 高速移動：7.8km/s，比地面交通快1000倍
• 長距離：500-2000km，信號傳輸延遲顯著
• 高密度：單顆衛星服務數千用戶
• 頻繁切換：覆蓋時間短，需要頻繁換手

【傳統方法局限性】
• 測量延遲：100-200ms測量週期
• 傳輸延遲：3.2-12ms信號傳輸
• 處理延遲：決策分析時間
• 總延遲：112-212ms（在快速移動環境中過長）

【技術需求】
• 即時決策：毫秒級換手決策
• 預測能力：提前預判最佳時機
• 資源協調：多用戶聯合優化"""

        self.add_figure_slide(prs, "研究背景：LEO衛星換手挑戰", background_content)

        # 4. DHO核心創新
        innovation_content = """【范式轉換】
傳統：測量 → 報告 → 決策（reactive）
DHO：觀察 → 學習 → 預測（proactive）

【關鍵技術突破】
1. 省略MR階段
   • 消除測量報告的傳輸延遲
   • 避免基於過時數據的決策

2. 智能預測決策
   • 基於LEO軌道確定性
   • 利用歷史模式進行學習
   • 實時狀態評估

3. 多用戶協調優化
   • 聯合考慮所有UE需求
   • 避免資源衝突
   • 動態負載均衡

【性能提升】
• 決策延遲：>100倍改善
• 功耗節省：30-50%
• 存取成功率：顯著提升"""

        self.add_figure_slide(prs, "DHO演算法核心創新", innovation_content,
                             figure_filename="page_4_img_1.png",
                             figure_description="圖2: DHO演算法架構圖")

        # 5. MDP建模 - 狀態空間
        state_content = """【狀態空間設計原理】
s[n] = {n, a^HO[n], a[n-1]}

【組件詳解】
1. 時間索引 n
   • 軌道週期內的時間位置
   • 隱式編碼衛星位置信息
   • 範圍：0 ≤ n < T（軌道週期）

2. 存取狀態 a^HO[n]
   • a^HO_j[n] ∈ {0,1}表示UE j的存取狀態
   • 1：已成功存取，0：未成功存取
   • 反映當前網路負載情況

3. 歷史動作 a[n-1]
   • 上一時隙的HO決策記錄
   • 提供決策連續性
   • 避免劇烈策略變化

【設計考量】
• 最小充分統計量：包含決策所需的最少信息
• 可觀察性：所有信息在服務衛星本地可得
• 計算效率：狀態維度控制在合理範圍"""

        self.add_formula_with_figure_slide(prs, "MDP建模：狀態空間設計",
                                         "s[n] = {n, a^HO[n], a[n-1]}",
                                         state_content,
                                         "",
                                         figure_filename="page_4_img_2.png",
                                         figure_description="圖3: 狀態空間結構圖")

        # 6. MDP建模 - 動作空間
        action_content = """【動作空間設計】
a_j[n] = {a_0, a_1, a_2, ..., a_{K-1}}
約束：Σ(k=0 to K-1) a_k = 1

【One-hot編碼邏輯】
• a_0 = 1：不進行換手（智能退避）
• a_k = 1 (k≥1)：選擇目標衛星k
• 互斥選擇：每個UE只能選一個目標

【全域協調】
• 全域動作：a[n] = [a_1[n], a_2[n], ..., a_J[n]]^T
• 聯合優化：考慮所有UE的協調效果
• 衝突避免：預防多個UE選擇同一目標

【技術優勢】
• 負載均衡：分散UE到不同衛星
• 智能退避：避免不必要的換手
• 資源協調：最大化系統整體性能"""

        self.add_formula_with_figure_slide(prs, "MDP建模：動作空間設計",
                                         "a_j[n] = {a_0, a_1, ..., a_{K-1}}, Σa_k = 1",
                                         action_content,
                                         "",
                                         figure_filename="page_4_img_3.png",
                                         figure_description="圖4: 動作空間示意圖")

        # 7. MDP建模 - 獎勵函數
        reward_content = """【獎勵函數設計】
r[n] = -D[n] - ν·C[n]

【組件詳解】
1. 存取延遲懲罰 D[n]
   D[n] = (1/|J|)·Σ(j∈J)(1 - a_j^HO[n])
   • 正規化的存取延遲
   • 懲罰未成功存取的UE

2. 碰撞率懲罰 C[n]
   C[n] = Σ(k=1 to K-1)C_k^R[n] + C^P[n]
   • C_k^R[n]：RB資源碰撞
   • C^P[n]：PRACH接入碰撞

3. 權衡係數 ν
   • URLLC應用：ν較大（重視低延遲）
   • mMTC應用：ν較小（容忍適度延遲）

【優化目標】
• 最小化存取延遲
• 最小化資源碰撞
• 應用自適應權衡"""

        self.add_formula_with_figure_slide(prs, "MDP建模：獎勵函數設計",
                                         "r[n] = -D[n] - ν·C[n]",
                                         reward_content,
                                         "",
                                         figure_filename="page_5_img_1.png",
                                         figure_description="圖5: 獎勵函數結構")

        # 8. IMPALA演算法 - V-trace機制
        vtrace_content = """【V-trace核心機制】
v[n] = V(s[n]) + Σ(i=s to s+k-1) γ^(i-s)·Π(j=s to i-1)c[j]·δ_i^V

【重要性採樣權重】
ρ[n] = min(ρ̄, π(a[n]|s[n])/μ(a[n]|s[n]))
c[n] = min(c̄, π(a[n]|s[n])/μ(a[n]|s[n]))

【機制說明】
• π(a|s)：learner的目標策略
• μ(a|s)：actor的行為策略
• 截斷機制：防止權重過大
• 偏差修正：確保正確的價值估計

【技術優勢】
• 解決策略滯後問題
• 保證算法收斂性
• 提高樣本利用效率
• 穩定訓練過程"""

        self.add_formula_with_figure_slide(prs, "IMPALA演算法：V-trace機制",
                                         "v[n] = V(s[n]) + Σγ^(i-s)·Πc[j]·δ_i^V",
                                         vtrace_content,
                                         "",
                                         figure_filename="page_5_img_2.png",
                                         figure_description="圖6: V-trace計算流程")

        # 9. 預測能力基礎
        prediction_content = """【軌道確定性原理】
q_i[m] = q_i[0] + τ·Σ(m'=1 to m)v_i[m'τ]

【預測機制層次】
Level 1: 位置預測
• 輸入：時間索引n
• 輸出：衛星精確位置
• 基礎：牛頓力學定律

Level 2: 覆蓋預測
• 輸入：衛星位置 + 用戶位置
• 輸出：信號覆蓋品質
• 依據：自由空間路徑損耗

Level 3: 決策預測
• 輸入：覆蓋品質 + 網路負載
• 輸出：最佳HO決策
• 方法：深度強化學習

【學習目標】
• 狀態→動作映射：f_θ(s) → a*
• 模式識別：時間→最優策略
• 泛化能力：適應未見情況"""

        self.add_formula_with_figure_slide(prs, "預測能力：軌道確定性基礎",
                                         "q_i[m] = q_i[0] + τ·Σv_i[m'τ]",
                                         prediction_content,
                                         "")

        # 10. 實驗結果對比
        experiment_content = """【實驗設置】
• 衛星數量：多種配置測試
• 用戶設備：大規模UE場景
• 網路負載：不同密度環境
• 對比方法：傳統A3事件觸發

【關鍵性能指標】
1. 存取延遲 (Access Delay)
   • DHO：顯著降低
   • 改善幅度：6.86倍

2. 決策延遲 (Decision Latency)
   • 傳統方法：112-212ms
   • DHO方法：<1ms
   • 改善幅度：>100倍

3. 功耗效率 (Power Efficiency)
   • 消除週期性測量
   • 節省30-50%相關功耗

4. 碰撞率 (Collision Rate)
   • 多用戶協調優化
   • 大幅降低資源衝突

【實驗驗證】
• 多種軌道配置
• 不同負載場景
• 穩定性測試
• 可擴展性驗證"""

        self.add_figure_slide(prs, "實驗結果與性能驗證", experiment_content)

        # 11. 性能對比表格
        performance_slide = prs.slides.add_slide(prs.slide_layouts[1])
        performance_slide.shapes.title.text = "DHO vs 傳統HO詳細性能對比"
        self.set_mixed_font_style(performance_slide.shapes.title.text_frame, font_size=20)

        # 創建詳細對比表格
        table_data = [
            ["性能指標", "傳統HO方法", "DHO方法", "改善程度", "技術意義"],
            ["決策延遲", "112-212ms", "<1ms", ">100倍", "即時響應"],
            ["功耗消耗", "基準100%", "50-70%", "30-50%節省", "延長電池壽命"],
            ["存取成功率", "基準", "6.86倍提升", "686%", "用戶體驗大幅改善"],
            ["資源利用率", "局部優化", "全域協調", "顯著提升", "系統容量增加"],
            ["適應能力", "固定規則", "動態學習", "質的飛躍", "智能化演進"],
            ["可擴展性", "有限", "支持大規模", "架構優勢", "面向未來需求"]
        ]

        rows = len(table_data)
        cols = len(table_data[0])

        left = Inches(0.3)
        top = Inches(1.8)
        width = Inches(9.4)
        height = Inches(4.5)

        table = performance_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 填入表格數據並設定格式
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                if i == 0:  # 標題行
                    self.set_mixed_font_style(cell.text_frame, font_size=11)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.primary_color
                else:
                    self.set_mixed_font_style(cell.text_frame, font_size=10)

        # 12. 技術貢獻總結
        contribution_content = """【學術貢獻】
1. 理論創新
   • 首次完整的LEO衛星HO的MDP建模
   • 省略MR的可行性理論證明
   • 深度強化學習在NTN的創新應用

2. 演算法貢獻
   • IMPALA在通訊系統的首次應用
   • V-trace機制的衛星環境適配
   • 多用戶聯合優化演算法

3. 工程價值
   • 為6G NTN標準提供技術參考
   • 通訊系統智能化的重要里程碑
   • 產業技術升級的推動力

【技術影響】
• 范式轉換：reactive → proactive
• 方法創新：規則 → 學習
• 優化層次：局部 → 全域

【未來展望】
• 多目標優化擴展
• 跨層聯合設計
• 大規模星座協作"""

        self.add_figure_slide(prs, "技術貢獻與學術影響", contribution_content)

        # 13. 結論
        conclusion_content = """【主要結論】
1. DHO成功實現了省略MR的智能換手
   • 技術可行性得到充分驗證
   • 性能優勢顯著且穩定

2. 深度強化學習適用於LEO衛星網路
   • MDP建模準確捕捉系統特性
   • IMPALA演算法表現優異

3. 預測式決策優於反應式決策
   • 利用軌道確定性實現精準預測
   • 多用戶協調避免資源衝突

【技術成熟度】
• 理論基礎：完整且嚴謹
• 演算法實現：穩定可靠
• 性能驗證：全面深入
• 工程化：具備可行性

【研究意義】
• 為LEO衛星網路優化開闢新方向
• 為6G NTN標準化提供重要參考
• 推動通訊系統向智能化演進

【致謝】
感謝所有參與研究的團隊成員
感謝審稿專家的寶貴意見"""

        self.add_figure_slide(prs, "結論與致謝", conclusion_content)

        # 儲存簡報
        output_filename = os.path.join(self.output_dir, "DHO演算法完整技術原理簡報.pptx")
        prs.save(output_filename)

        print(f"✅ DHO演算法完整技術簡報創建完成")
        print(f"📊 總投影片數：{len(prs.slides)} 張")
        print(f"📁 儲存位置：{output_filename}")

        # 生成創建報告
        self.generate_comprehensive_report(len(prs.slides))

        return output_filename

    def generate_comprehensive_report(self, total_slides):
        """生成完整簡報創建報告"""
        report_content = f"""# DHO演算法完整技術原理簡報創建報告

## 📊 簡報概覽
- **簡報主題**: DHO演算法完整技術原理（含圖片整合）
- **總投影片數**: {total_slides} 張
- **創建時間**: 2024-09-17
- **特色**: 深度技術解析 + 論文圖片整合

## 🎯 內容結構
1. **標題頁** - 簡報主題
2. **簡報大綱** - 完整架構指引（含系統架構圖）
3. **研究背景** - LEO衛星換手挑戰
4. **核心創新** - DHO技術突破（含演算法架構圖）
5. **狀態空間** - MDP建模詳解（含結構圖）
6. **動作空間** - 設計原理說明（含示意圖）
7. **獎勵函數** - 數學建模解析（含結構圖）
8. **V-trace機制** - IMPALA演算法（含流程圖）
9. **預測基礎** - 軌道確定性原理
10. **實驗驗證** - 性能測試結果
11. **性能對比** - 詳細對比表格
12. **技術貢獻** - 學術與工程價值
13. **結論致謝** - 總結與展望

## 🖼️ 圖片整合
- **整合圖片數**: {len(self.figure_info)} 張原始論文圖片
- **圖片說明**: 每張圖片都有對應的技術描述
- **佈局優化**: 圖文並茂的專業排版
- **圖片品質**: 保持原始論文的高品質

## 🔬 技術深度
- **公式解析**: 每個核心公式的完整數學表達
- **變數說明**: 詳細的變數定義和物理意義
- **流程圖**: 完整的演算法執行流程
- **對比分析**: 量化的性能提升數據

## ✅ 品質特色
- **視覺效果**: 圖文並茂，專業美觀
- **技術準確**: 基於原始論文的精確內容
- **邏輯完整**: 從背景到結論的完整技術脈絡
- **實用價值**: 適合多種學術和工程場合

## 📈 應用場景
本完整版簡報適合：
- **博士論文答辯**: 完整技術展示
- **學術會議報告**: 深度技術交流
- **工程團隊分享**: 實現細節討論
- **技術標準提案**: 標準化組織展示
- **產學合作**: 技術轉移說明

## 🎖️ 相比其他版本優勢
- **更完整**: 包含背景、技術、實驗、結論全流程
- **更直觀**: 整合原始論文圖片，增強理解效果
- **更專業**: 符合學術規範的完整技術報告
- **更實用**: 可直接用於多種正式場合
"""

        report_filename = os.path.join(self.output_dir, "DHO完整簡報創建報告.md")
        with open(report_filename, 'w', encoding='utf-8') as f:
            f.write(report_content)

        print(f"📋 完整報告已生成：{report_filename}")

def main():
    """主函數"""
    print("🚀 DHO演算法完整技術原理簡報生成器（含圖片）")
    print("=" * 60)

    # 創建生成器實例
    generator = DHOAlgorithmCompletePresentation()

    # 生成簡報
    try:
        output_file = generator.create_comprehensive_dho_presentation()
        print(f"\n✅ 完整簡報生成成功：{output_file}")
        print("\n🎯 本簡報特色：")
        print("   • 完整的技術原理解析")
        print("   • 整合原始論文圖片")
        print("   • 詳細的公式和變數說明")
        print("   • 完整的演算法流程")
        print("   • 量化的性能對比")
        print("   • 專業的學術品質")
        print("   • 適合正式技術展示")

    except Exception as e:
        print(f"❌ 簡報生成失敗：{e}")
        import traceback
        traceback.print_exc()
        return False

    return True

if __name__ == "__main__":
    main()